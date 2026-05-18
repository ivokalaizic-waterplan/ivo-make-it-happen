#!/usr/bin/env python3
"""Generate Q1 2026 Board Meeting slides — Blue Team (Ivo's scope).
   Slide 1: Target Tracking Q1 (split layout — text left, screenshot placeholder right)
   Slide 2: Platform Navigation Q1 (split layout)
   Slide 3: Q2 Roadmap (card layout — no screenshots)
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Palette ────────────────────────────────────────────────────────────
FIREFLY_500 = RGBColor(0x35, 0xA1, 0xC9)
FIREFLY_700 = RGBColor(0x20, 0x60, 0x78)
FIREFLY_900 = RGBColor(0x0F, 0x2D, 0x38)
FIREFLY_100 = RGBColor(0xD6, 0xEC, 0xF4)
FIREFLY_50  = RGBColor(0xEB, 0xF6, 0xFA)
GRAY_900    = RGBColor(0x2A, 0x2B, 0x2D)
GRAY_700    = RGBColor(0x55, 0x57, 0x5B)
GRAY_500    = RGBColor(0x80, 0x82, 0x86)
GRAY_300    = RGBColor(0xB0, 0xB2, 0xB5)
GRAY_100    = RGBColor(0xDA, 0xDB, 0xDE)
GRAY_50     = RGBColor(0xEE, 0xEF, 0xF0)
WHITE       = RGBColor(0xFD, 0xFD, 0xFD)
SLIDE_BG    = RGBColor(0xF3, 0xF5, 0xF8)
SECTION_BG  = RGBColor(0xEC, 0xEE, 0xF1)
CARD_BG     = RGBColor(0xFF, 0xFF, 0xFF)
SUCCESS     = RGBColor(0x02, 0x7A, 0x48)

SKILL_DIR = os.path.expanduser("~/.claude/skills/waterplan-pptx")
ASSETS    = os.path.join(SKILL_DIR, "assets")

SW = Inches(13.333)
SH = Inches(7.5)

# ── Helpers ────────────────────────────────────────────────────────────

def set_slide_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_accent_bar(slide, top, color=FIREFLY_500, height=Pt(4)):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), top, SW, height)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()


def add_waterplan_logo(slide):
    logo = os.path.join(ASSETS, "waterplan-logo-dark.png")
    slide.shapes.add_picture(logo, SW - Inches(2.2), Inches(0.4), Inches(1.5), Inches(0.30))


def add_confidential(slide):
    tb = slide.shapes.add_textbox(SW - Inches(2.0), SH - Inches(0.45), Inches(1.5), Inches(0.25))
    p = tb.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    r = p.add_run()
    r.text = "CONFIDENTIAL"
    r.font.size = Pt(8)
    r.font.color.rgb = GRAY_500
    r.font.name = "Inter"
    r.font.bold = True


def slide_chrome(slide):
    set_slide_bg(slide, SLIDE_BG)
    add_accent_bar(slide, Inches(0))
    add_accent_bar(slide, SH - Pt(4))
    add_waterplan_logo(slide)
    add_confidential(slide)


def add_title(slide, title, subtitle=None):
    tb = slide.shapes.add_textbox(Inches(0.7), Inches(0.35), Inches(8), Inches(0.55))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.size = Pt(26)
    r.font.color.rgb = FIREFLY_700
    r.font.name = "Inter"
    r.font.bold = True
    if subtitle:
        r2 = p.add_run()
        r2.text = f"\n{subtitle}"
        r2.font.size = Pt(14)
        r2.font.color.rgb = GRAY_500
        r2.font.name = "Inter"
        r2.font.bold = False


def add_section_label(slide, left, top, text, color=FIREFLY_500):
    tb = slide.shapes.add_textbox(left, top, Inches(4), Inches(0.3))
    p = tb.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = text
    r.font.size = Pt(14)
    r.font.color.rgb = color
    r.font.name = "Inter"
    r.font.bold = True


def add_body(slide, left, top, width, text, font_size=13, color=GRAY_700):
    tb = slide.shapes.add_textbox(left, top, width, Inches(0.01))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.line_spacing = 1.4
    r = p.add_run()
    r.text = text
    r.font.size = Pt(font_size)
    r.font.color.rgb = color
    r.font.name = "Inter"
    return tb


def add_bullets(slide, left, top, width, bullets, font_size=13, color=GRAY_900):
    """bullets = list of str"""
    tb = slide.shapes.add_textbox(left, top, width, Inches(0.01))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = 1.5
        p.space_after = Pt(4)
        r = p.add_run()
        r.text = f"•   {b}"
        r.font.size = Pt(font_size)
        r.font.color.rgb = color
        r.font.name = "Inter"
    return tb


def add_screenshot_placeholder(slide, left, top, width, height, label="Add screenshot"):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    s.fill.solid()
    s.fill.fore_color.rgb = RGBColor(0xF8, 0xF9, 0xFB)
    s.line.color.rgb = GRAY_100
    s.line.width = Pt(1.5)
    s.line.dash_style = 2
    s.adjustments[0] = 0.03
    tf = s.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = label
    r.font.size = Pt(11)
    r.font.color.rgb = GRAY_300
    r.font.name = "Inter"
    r.font.italic = True


def add_card(slide, left, top, width, height):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    s.fill.solid()
    s.fill.fore_color.rgb = CARD_BG
    s.line.color.rgb = GRAY_100
    s.line.width = Pt(0.75)
    s.adjustments[0] = 0.05
    return s


def add_pill(slide, left, top, width, height, fill, text, text_color=WHITE, font_size=11):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    s.line.fill.background()
    s.adjustments[0] = 0.5
    tf = s.text_frame
    tf.word_wrap = False
    tf.margin_left = Pt(8)
    tf.margin_right = Pt(8)
    tf.margin_top = Pt(0)
    tf.margin_bottom = Pt(0)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.size = Pt(font_size)
    r.font.color.rgb = text_color
    r.font.name = "Inter"
    r.font.bold = True


def add_rich_bullets(slide, left, top, width, bullets, font_size=12,
                     color=GRAY_900, bold_color=FIREFLY_900):
    """bullets = list of list-of-(text, bold) tuples."""
    tb = slide.shapes.add_textbox(left, top, width, Inches(0.01))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, runs_data in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = 1.45
        p.space_after = Pt(6)
        rb = p.add_run()
        rb.text = "•   "
        rb.font.size = Pt(font_size)
        rb.font.color.rgb = color
        rb.font.name = "Inter"
        for txt, is_bold in runs_data:
            r = p.add_run()
            r.text = txt
            r.font.size = Pt(font_size)
            r.font.color.rgb = bold_color if is_bold else color
            r.font.name = "Inter"
            r.font.bold = is_bold
    return tb


# ── Presentation ───────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width = SW
prs.slide_height = SH
blank = prs.slide_layouts[6]

# Text column boundaries (for split layout)
TXT_L = Inches(0.7)
TXT_W = Inches(5.3)

# Screenshot area (right side)
IMG_L = Inches(6.4)
IMG_W = Inches(6.2)


# ======================================================================
# SLIDE 1 — Target Tracking Q1 (split layout)
# ======================================================================
s1 = prs.slides.add_slide(blank)
slide_chrome(s1)

add_title(s1, "Target Tracking", "Q1 2026 Progress")

# Subtitle line
tb = s1.shapes.add_textbox(TXT_L, Inches(1.05), TXT_W, Inches(0.3))
p = tb.text_frame.paragraphs[0]
r = p.add_run()
r.text = "Yearly Targets, Governance & Configurable Views"
r.font.size = Pt(15)
r.font.color.rgb = FIREFLY_500
r.font.name = "Inter"
r.font.italic = True
r.font.bold = False

# What it aims to achieve
add_section_label(s1, TXT_L, Inches(1.55), "What it aims to achieve")
add_bullets(s1, TXT_L, Inches(1.9), TXT_W, [
    "Enable companies to define year-by-year target milestones.",
    "Protect committed values from unauthorized edits.",
    "Let each company configure which KPIs, charts, and filters are visible in their Target Tracking experience.",
], font_size=13)

# Why it matters
add_section_label(s1, TXT_L, Inches(3.55), "Why it matters for our users")
add_bullets(s1, TXT_L, Inches(3.9), TXT_W, [
    "Yearly targets enable earlier deviation detection, credible roadmaps, and site-level planning.",
    "Locking ensures audit readiness.",
    "Configurable views reduce noise and increase adoption.",
], font_size=13)

# Screenshot placeholders (right side)
add_screenshot_placeholder(s1, IMG_L, Inches(0.9), IMG_W, Inches(3.0),
                           "Screenshot: Target evolution chart with yearly checkpoints")
add_screenshot_placeholder(s1, IMG_L + Inches(0.5), Inches(4.2), Inches(2.6), Inches(2.7),
                           "Screenshot: Edit target + Lock")
add_screenshot_placeholder(s1, IMG_L + Inches(3.3), Inches(4.2), Inches(2.6), Inches(2.7),
                           "Screenshot: Target tracking settings")


# ======================================================================
# SLIDE 2 — Platform Navigation & Enterprise Customization Q1
# ======================================================================
s2 = prs.slides.add_slide(blank)
slide_chrome(s2)

add_title(s2, "Platform Navigation & Enterprise Customization", "Q1 2026 Progress")

# What it aims to achieve
add_section_label(s2, TXT_L, Inches(1.35), "What it aims to achieve")
add_bullets(s2, TXT_L, Inches(1.7), TXT_W, [
    "Replace the legacy horizontal navigation with a scalable side navigation.",
    "Allow companies to rename platform features to match their internal terminology.",
    "Build a configurable homepage so users land directly on what matters to them — scaling from a first live deployment to a fully configurable system.",
], font_size=13)

# Why it matters
add_section_label(s2, TXT_L, Inches(3.35), "Why it matters for our users")
add_bullets(s2, TXT_L, Inches(3.7), TXT_W, [
    "As we add more modules, the old horizontal navigation breaks down.",
    "Enterprise clients need the platform to speak their language.",
], font_size=13)

add_body(s2, TXT_L, Inches(4.65), TXT_W,
         "These changes reduce friction, improve time-to-value, and make the platform feel native to each client.",
         font_size=13, color=GRAY_700)

# Rollout status line
add_body(s2, TXT_L, Inches(5.35), TXT_W,
         "Side navigation and homepage currently live for internal users and select clients. Full rollout planned for Q2.",
         font_size=11, color=GRAY_500)

# Screenshot placeholders (right side)
add_screenshot_placeholder(s2, IMG_L, Inches(0.9), Inches(2.8), Inches(5.5),
                           "Screenshot: Side navigation")
add_screenshot_placeholder(s2, IMG_L + Inches(3.1), Inches(0.9), Inches(3.1), Inches(5.5),
                           "Screenshot: Homepage + Custom module names")


# ======================================================================
# SLIDE 3 — Q2 2026 Roadmap (card layout)
# ======================================================================
s3 = prs.slides.add_slide(blank)
slide_chrome(s3)

add_title(s3, "Q2 Roadmap — Target Tracking, Platform Experience & Enterprise Scale")

# ── Three initiative cards ──
card_w = Inches(3.7)
card_h = Inches(4.5)
card_top = Inches(1.35)
gap = Inches(0.3)
cards_total = card_w * 3 + gap * 2
start_left = (SW - cards_total) / 2

cards = [
    {
        "pill": "AI Scenario Prioritization",
        "pill_w": Inches(2.8),
        "pill_color": FIREFLY_500,
        "what_title": "What it aims to achieve",
        "what_bullets": [
            "AI-powered project prioritization given budget, operational constraints, and a reduction target.",
            "The platform recommends which projects to execute and in what order — replacing manual multi-criteria analysis.",
        ],
        "why_title": "Why it matters for our users",
        "why_bullets": [
            [("Committed deliverable for AB InBev", True), (" and a major product differentiator.", False)],
            [("Turns Waterplan into a decision engine", True), (" — from tracking progress to recommending action.", False)],
            [("Gives managers a defensible, data-backed recommendation", True), (" they can present to leadership for budget approval.", False)],
        ],
    },
    {
        "pill": "Scale Platform Experience",
        "pill_w": Inches(2.7),
        "pill_color": FIREFLY_500,
        "what_title": "What it aims to achieve",
        "what_bullets": [
            "Roll out side navigation and configurable homepage to all clients — building on the successful first deployment for Colgate.",
            "Develop full mobile responsive experience to enable field teams on mobile devices.",
        ],
        "why_title": "Why it matters for our users",
        "why_bullets": [
            [("Every client gets a personalized entry point", True), (" — users land on what's relevant to their role, not a generic overview.", False)],
            [("Mobile support unlocks AI Water Meters", True), (" for plant operators who work from their phones, not desktops.", False)],
            [("Scaling from one client to all", True), (" validates the architecture and accelerates adoption across the portfolio.", False)],
        ],
    },
    {
        "pill": "Enterprise Access Control",
        "pill_w": Inches(2.6),
        "pill_color": FIREFLY_500,
        "what_title": "What it aims to achieve",
        "what_bullets": [
            "Role-based action permissions — Company Owners define what each role can do: create targets, edit, lock.",
            "Configurable cross-site visibility — match the platform to each client's corporate data access policies.",
        ],
        "why_title": "Why it matters for our users",
        "why_bullets": [
            [("Without governance, companies can't safely expand", True), (" Target Tracking to more users and sites.", False)],
            [("Coca-Cola needs all-sites-visible", True), (" for coordination; ", False), ("Colgate needs strict restriction", True), (" for data confidentiality.", False)],
            [("Removes the last blocker", True), (" for large-scale enterprise rollout across hundreds of users.", False)],
        ],
    },
]

for i, c in enumerate(cards):
    left = start_left + i * (card_w + gap)

    add_card(s3, left, card_top, card_w, card_h)

    # Pill
    add_pill(s3, left + Inches(0.25), card_top + Inches(0.2),
             c["pill_w"], Inches(0.3), c["pill_color"], c["pill"])

    # What it aims to achieve
    add_section_label(s3, left + Inches(0.25), card_top + Inches(0.65),
                      c["what_title"], FIREFLY_500)
    add_bullets(s3, left + Inches(0.25), card_top + Inches(0.95),
                card_w - Inches(0.5), c["what_bullets"], font_size=11, color=GRAY_900)

    # Why it matters
    why_top = card_top + Inches(2.35)
    add_section_label(s3, left + Inches(0.25), why_top,
                      c["why_title"], FIREFLY_500)
    add_rich_bullets(s3, left + Inches(0.25), why_top + Inches(0.3),
                     card_w - Inches(0.5), c["why_bullets"],
                     font_size=11, color=GRAY_900, bold_color=FIREFLY_900)


# ── Scope 3 footer line ──
tb = s3.shapes.add_textbox(start_left, card_top + card_h + Inches(0.25),
                           cards_total, Inches(0.35))
tf = tb.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.LEFT

r1 = p.add_run()
r1.text = "Also exploring:  "
r1.font.size = Pt(11)
r1.font.color.rgb = GRAY_500
r1.font.name = "Inter"
r1.font.bold = False

r2 = p.add_run()
r2.text = "Scope 3 Logistics"
r2.font.size = Pt(11)
r2.font.color.rgb = FIREFLY_700
r2.font.name = "Inter"
r2.font.bold = True

r3 = p.add_run()
r3.text = " — initiating discovery for Scope 3 logistics emissions tracking, expanding carbon capabilities beyond Scope 1 & 2."
r3.font.size = Pt(11)
r3.font.color.rgb = GRAY_500
r3.font.name = "Inter"
r3.font.bold = False


# ── Save ───────────────────────────────────────────────────────────────
output_path = os.path.expanduser("~/Downloads/Board-Q1-2026-Blue-Team-Ivo.pptx")
prs.save(output_path)
print(f"Saved → {output_path}")
