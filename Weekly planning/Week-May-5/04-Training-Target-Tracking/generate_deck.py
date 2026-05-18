#!/usr/bin/env python3
"""Generate Target Tracking Training deck for Waterplan — May 2026."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import os

# ── Palette (Seastem / Firefly — internal deck) ──────────────────────────
FIREFLY_500 = RGBColor(0x35, 0xA1, 0xC9)
FIREFLY_700 = RGBColor(0x20, 0x60, 0x78)
FIREFLY_900 = RGBColor(0x0F, 0x2D, 0x38)
FIREFLY_100 = RGBColor(0xD6, 0xEC, 0xF4)
FIREFLY_50  = RGBColor(0xEB, 0xF6, 0xFA)
GRAY_900    = RGBColor(0x2A, 0x2B, 0x2D)
GRAY_700    = RGBColor(0x55, 0x57, 0x5B)
GRAY_100    = RGBColor(0xDA, 0xDB, 0xDE)
GRAY_50     = RGBColor(0xEE, 0xEF, 0xF0)
WHITE       = RGBColor(0xFD, 0xFD, 0xFD)
SUCCESS     = RGBColor(0x02, 0x7A, 0x48)
WARNING     = RGBColor(0xB5, 0x47, 0x08)
ERROR       = RGBColor(0xB3, 0x23, 0x18)

FONT = "Inter"

# ── Helpers ───────────────────────────────────────────────────────────────

def add_waterplan_logo(slide):
    logo_left = Inches(10.8)
    logo_top = Inches(0.3)
    logo_w = Inches(2)
    logo_h = Inches(0.45)
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, logo_left, logo_top, logo_w, logo_h)
    shape.fill.background()
    shape.line.color.rgb = FIREFLY_500
    shape.line.width = Pt(1.5)
    shape.adjustments[0] = 0.5
    tf = shape.text_frame
    tf.margin_left = Pt(8); tf.margin_right = Pt(8)
    tf.margin_top = Pt(0); tf.margin_bottom = Pt(0)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "waterplan"
    run.font.size = Pt(13)
    run.font.color.rgb = FIREFLY_500
    run.font.name = FONT
    run.font.bold = False


def add_title(slide, text, highlight_word=None, y=Inches(0.3), font_size=Pt(24)):
    txBox = slide.shapes.add_textbox(Inches(0.5), y, Inches(10), Inches(0.9))
    tf = txBox.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    if highlight_word and highlight_word in text:
        parts = text.split(highlight_word, 1)
        for i, (txt, color) in enumerate([(parts[0], FIREFLY_700), (highlight_word, FIREFLY_500), (parts[1], FIREFLY_700)]):
            if txt:
                r = p.add_run(); r.text = txt; r.font.size = font_size
                r.font.color.rgb = color; r.font.name = FONT; r.font.bold = True
    else:
        r = p.add_run(); r.text = text; r.font.size = font_size
        r.font.color.rgb = FIREFLY_700; r.font.name = FONT; r.font.bold = True


def add_bullets(slide, items, left=Inches(0.7), top=Inches(1.5), width=Inches(11.5), font_size=Pt(14), spacing=Pt(8), color=GRAY_900):
    txBox = slide.shapes.add_textbox(left, top, width, Inches(5))
    tf = txBox.text_frame; tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = spacing
        p.level = 0
        r = p.add_run()
        r.text = f"\u2022  {item}"
        r.font.size = font_size; r.font.color.rgb = color; r.font.name = FONT


def add_note(slide, text, top=Inches(6.0), left=Inches(0.7)):
    txBox = slide.shapes.add_textbox(left, top, Inches(11), Inches(0.5))
    tf = txBox.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = text
    r.font.size = Pt(11); r.font.color.rgb = GRAY_700; r.font.name = FONT; r.font.italic = True


def add_screenshot_placeholder(slide, label, left=Inches(0.7), top=Inches(4.0), width=Inches(11.5), height=Inches(3.0)):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = GRAY_50
    shape.line.color.rgb = GRAY_100
    shape.line.width = Pt(1.5)
    shape.line.dash_style = 2  # dash
    shape.adjustments[0] = 0.05
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].space_before = Pt(0)
    # vertical center
    txBody = tf._txBody
    bodyPr = txBody.find(qn('a:bodyPr'))
    bodyPr.set('anchor', 'ctr')
    r = p.add_run()
    r.text = f"[SCREENSHOT: {label}]"
    r.font.size = Pt(16); r.font.color.rgb = GRAY_700; r.font.name = FONT; r.font.bold = True


def add_styled_table(slide, rows_data, left, top, width, height):
    n_rows = len(rows_data); n_cols = len(rows_data[0])
    ts = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
    table = ts.table
    # Remove default table style banding
    tbl = ts._element.graphic.graphicData.tbl
    tblPr = tbl.find(qn('a:tblPr'))
    if tblPr is not None:
        tblPr.set('bandRow', '0')

    for ri, row in enumerate(rows_data):
        for ci, cell_text in enumerate(row):
            cell = table.cell(ri, ci)
            cell.text = ""
            tf = cell.text_frame; tf.word_wrap = True
            p = tf.paragraphs[0]
            r = p.add_run(); r.text = str(cell_text)
            r.font.size = Pt(11); r.font.name = FONT
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            tf.margin_left = Pt(6); tf.margin_right = Pt(6)
            tf.margin_top = Pt(4); tf.margin_bottom = Pt(4)
            if ri == 0:
                cell.fill.solid(); cell.fill.fore_color.rgb = FIREFLY_700
                r.font.color.rgb = WHITE; r.font.bold = True
            elif ri % 2 == 0:
                cell.fill.solid(); cell.fill.fore_color.rgb = FIREFLY_50
                r.font.color.rgb = GRAY_900
            else:
                cell.fill.solid(); cell.fill.fore_color.rgb = WHITE
                r.font.color.rgb = GRAY_900
    return ts


def add_section_label(slide, text, top=Inches(1.15)):
    txBox = slide.shapes.add_textbox(Inches(0.5), top, Inches(5), Inches(0.35))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = text.upper()
    r.font.size = Pt(10); r.font.color.rgb = FIREFLY_500; r.font.name = FONT; r.font.bold = True


def new_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    # white background
    bg = slide.background
    fill = bg.fill; fill.solid(); fill.fore_color.rgb = WHITE
    add_waterplan_logo(slide)
    return slide


# ── Build presentation ───────────────────────────────────────────────────
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Cover
# ═══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
bg = s.background; fill = bg.fill; fill.solid(); fill.fore_color.rgb = FIREFLY_900

# Logo on cover
logo_shape = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.5), Inches(0.6), Inches(2.3), Inches(0.5))
logo_shape.fill.background()
logo_shape.line.color.rgb = FIREFLY_500; logo_shape.line.width = Pt(1.5)
logo_shape.adjustments[0] = 0.5
tf = logo_shape.text_frame
tf.margin_left = Pt(8); tf.margin_right = Pt(8); tf.margin_top = Pt(0); tf.margin_bottom = Pt(0)
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = "waterplan"; r.font.size = Pt(14); r.font.color.rgb = FIREFLY_500; r.font.name = FONT

# Title
txBox = s.shapes.add_textbox(Inches(1.5), Inches(2.3), Inches(10.3), Inches(1.5))
tf = txBox.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = "Target Tracking Training"
r.font.size = Pt(40); r.font.color.rgb = WHITE; r.font.name = FONT; r.font.bold = True

# Subtitle
txBox2 = s.shapes.add_textbox(Inches(1.5), Inches(3.8), Inches(10.3), Inches(0.8))
tf2 = txBox2.text_frame; tf2.word_wrap = True
p2 = tf2.paragraphs[0]; p2.alignment = PP_ALIGN.CENTER
r2 = p2.add_run(); r2.text = "For Solution Specialists & Sales"
r2.font.size = Pt(20); r2.font.color.rgb = FIREFLY_100; r2.font.name = FONT

# Bottom
txBox3 = s.shapes.add_textbox(Inches(1.5), Inches(5.8), Inches(10.3), Inches(0.6))
tf3 = txBox3.text_frame
p3 = tf3.paragraphs[0]; p3.alignment = PP_ALIGN.CENTER
r3 = p3.add_run(); r3.text = "May 2026 \u2014 Internal Training"
r3.font.size = Pt(14); r3.font.color.rgb = FIREFLY_500; r3.font.name = FONT

# Accent line under title
line = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.5), Inches(3.55), Inches(2.3), Pt(3))
line.fill.solid(); line.fill.fore_color.rgb = FIREFLY_500
line.line.fill.background()

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Why This Matters Now
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide(prs)
add_title(s, "The world changed \u2014 companies must track their environmental goals", highlight_word="track")
add_bullets(s, [
    "Regulations (CSRD, SEC, IFRS S2) now require reporting on targets and progress",
    "Investors demand measurable ESG targets with evidence of progress",
    "80% of companies still use spreadsheets to track their goals",
])

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 3 — How Companies Track Goals Today
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide(prs)
add_title(s, "The current state: spreadsheets and static reports", highlight_word="spreadsheets")
add_bullets(s, [
    "Spreadsheets: manual data, no connection to sites, error-prone",
    "BI tools (Domo, Tableau): separate pipelines, not designed for environmental management",
    "Static reports: outdated before they are finished",
    "Result: the sustainability manager spends 2 weeks preparing a quarterly review",
])

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 4 — The Real Pain Points
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide(prs)
add_title(s, "What we hear from clients again and again", highlight_word="clients")
add_bullets(s, [
    "\u201cI don\u2019t know if my reduction plan is working\u201d",
    "\u201cMy sites have different reduction timelines\u201d",
    "\u201cI can\u2019t show the board if our projects are enough\u201d",
    "\u201cThe quarterly review takes 2 weeks to prepare\u201d",
    "\u201cI need to compare different investment strategies\u201d",
], font_size=Pt(15))

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 5 — The Cost of Doing Nothing
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide(prs)
add_title(s, "What happens without structured tracking", highlight_word="without")
add_bullets(s, [
    "Credibility risk: published goals without evidence = greenwashing",
    "Money lost: without cost analysis, they invest in the wrong projects",
    "Missed reporting: can\u2019t respond to CDP, SBTi, CSRD with auditable data",
    "No accountability: corporate says \u2018reduce 30%\u2019 but sites have no individual target",
])

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Waterplan Target Tracking
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide(prs)
add_title(s, "One platform for water + carbon + energy", highlight_word="one platform")
add_section_label(s, "WATERPLAN TARGET TRACKING")
add_bullets(s, [
    "One target, three levels: Company \u2192 Site Group \u2192 Site",
    "Two domains: Water targets + Carbon & Energy targets (same engine, same UX)",
    "Five chart types: Evolution, Scenario Gap, MACC, Waterfall Impact, Waterfall Cost",
    "Up to 5 scenarios per target to compare investment strategies",
    "Automatic on-track / off-track based on forecast + projects",
    "Target locking for corporate governance",
], top=Inches(1.7))

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Spreadsheet vs Waterplan (table)
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide(prs)
add_title(s, "Side by side: what changes", highlight_word="what changes")
table_data = [
    ["What they do today", "What Waterplan TT does"],
    ["Targets disconnected from site data", "Targets connected to real-time operational data"],
    ["No scenario modeling", "Up to 5 scenarios with automatic impact calculation"],
    ["Manual on-track evaluation", "Automatic on-track/off-track updated with each data refresh"],
    ["Static charts rebuilt by hand", "5 interactive chart types, downloadable as PNG"],
    ["No audit trail", "SMART fields, activity log with attribution"],
    ["Single view (company or site)", "Company + Site + Site Group with independent targets"],
]
add_styled_table(s, table_data, Inches(0.7), Inches(1.5), Inches(11.8), Inches(5.0))

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Target Creation
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide(prs)
add_title(s, "One target, three levels of detail", highlight_word="three levels")
add_bullets(s, [
    "Company level: board commitments, CDP, SBTi goals",
    "Site level: operational goals per facility (local accountability)",
    "Site Group level: regional or business unit goals",
], top=Inches(1.5))
add_note(s, "Each level is independent \u2014 site targets don\u2019t need to equal the company target", top=Inches(3.5))
add_screenshot_placeholder(s, "Target creation modal", top=Inches(4.2), height=Inches(2.8))

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Evolution Chart
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide(prs)
add_title(s, "The chart that tells the story", highlight_word="story")
add_bullets(s, [
    "Blue line: historical data (what actually happened)",
    "Dashed line: forecast (what happens if we do nothing)",
    "Straight line: target (where we want to be)",
    "Scenario line: forecast adjusted with project impact",
], top=Inches(1.5))
add_screenshot_placeholder(s, "Evolution Chart with 4 lines", top=Inches(3.8), height=Inches(3.2))

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Scenario Modeling
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide(prs)
add_title(s, "Model before you invest \u2014 what-if scenarios", highlight_word="what-if scenarios")
add_bullets(s, [
    "Up to 5 scenarios per target (recommended: conservative / planned / stretch)",
    "Each scenario = a selection of projects from the pipeline",
    "The system automatically calculates: status, gap, total cost, cost per unit",
    "Result: on-track or off-track with exact gap percentage",
], top=Inches(1.5))
add_screenshot_placeholder(s, "Scenario selector with KPI cards", top=Inches(3.8), height=Inches(3.2))

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 11 — MACC Curve
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide(prs)
add_title(s, "MACC \u2014 Which projects give the most impact per dollar", highlight_word="MACC")
add_bullets(s, [
    "Available only for energy/carbon targets",
    "X axis: cumulative impact (emission reduction)",
    "Y axis: marginal cost per unit ($/tCO2e)",
    "Projects on the left = quick wins (do first). Projects on the right = expensive (do last)",
], top=Inches(1.5))
add_screenshot_placeholder(s, "MACC Curve", top=Inches(3.8), height=Inches(3.2))

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 12 — Waterfall Charts
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide(prs)
add_title(s, "Where is the impact and where is the cost", highlight_word="impact")
add_bullets(s, [
    "Impact Waterfall: from baseline \u2192 each project\u2019s contribution \u2192 projected vs target",
    "Cost Waterfall: cost per project (electricity in yellow, thermal in blue)",
    "Downloadable as PNG for presentations",
], top=Inches(1.5))
add_screenshot_placeholder(s, "Impact Waterfall", top=Inches(3.5), height=Inches(3.5))

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 13 — Target Locking
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide(prs)
add_title(s, "When the target is approved, it freezes", highlight_word="freezes")
add_bullets(s, [
    "Only Company Owners can lock/unlock",
    "A locked target is 100% read-only: no editing, no deleting",
    "Protects published goals (CDP, SBTi, annual reports)",
    "Activity log tracks who locked it, when, and any change attempts",
], top=Inches(1.5))
add_screenshot_placeholder(s, "Lock button in sidepanel", top=Inches(3.8), height=Inches(3.2))

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 14 — Alerts and Thresholds
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide(prs)
add_title(s, "Don\u2019t wait for the quarterly review", highlight_word="Don\u2019t wait")
add_bullets(s, [
    "Cluster alerts: when a site changes its target, corporate gets a notification",
    "Activity log shows the complete history of changes",
    "Thresholds: configurable limits per metric for proactive monitoring",
], top=Inches(1.5))
add_screenshot_placeholder(s, "Alert icon on target card", top=Inches(3.5), height=Inches(3.5))

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 15 — Company View: Aggregation
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide(prs)
add_title(s, "From macro to micro \u2014 company view with drill-down", highlight_word="drill-down")
add_bullets(s, [
    "Company view shows aggregated data from all sites",
    "\u2018Breakdown\u2019 button shows a table with each site and its individual target",
    "Detect which sites are on-track and which are dragging the average",
], top=Inches(1.5))
add_screenshot_placeholder(s, "Company view with Breakdown modal", top=Inches(3.5), height=Inches(3.5))

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 16 — Live Demo Transition
# ═══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
bg = s.background; fill = bg.fill; fill.solid(); fill.fore_color.rgb = FIREFLY_900
# Logo
add_waterplan_logo(s)
# override logo color for dark bg
for shape in s.shapes:
    pass  # logo is already light colored

txBox = s.shapes.add_textbox(Inches(1.5), Inches(2.0), Inches(10.3), Inches(1.2))
tf = txBox.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = "Let\u2019s see this in action"
r.font.size = Pt(36); r.font.color.rgb = WHITE; r.font.name = FONT; r.font.bold = True

# Bullets centered
items = [
    "Sales demo platform with real data",
    "9 sites in 9 countries",
    "Water targets + Carbon & Energy targets",
]
txBox2 = s.shapes.add_textbox(Inches(3), Inches(3.5), Inches(7.3), Inches(2.5))
tf2 = txBox2.text_frame; tf2.word_wrap = True
for i, item in enumerate(items):
    p2 = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER; p2.space_after = Pt(10)
    r2 = p2.add_run(); r2.text = f"\u2022  {item}"
    r2.font.size = Pt(16); r2.font.color.rgb = FIREFLY_100; r2.font.name = FONT

txBox3 = s.shapes.add_textbox(Inches(3), Inches(5.5), Inches(7.3), Inches(0.5))
tf3 = txBox3.text_frame
p3 = tf3.paragraphs[0]; p3.alignment = PP_ALIGN.CENTER
r3 = p3.add_run(); r3.text = "~35 minutes walkthrough"
r3.font.size = Pt(12); r3.font.color.rgb = FIREFLY_500; r3.font.name = FONT; r3.font.italic = True

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 17 — Scope 3: Coming Soon
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide(prs)
add_title(s, "Coming soon: Scope 3 \u2014 Value chain emissions", highlight_word="Scope 3")

# Three cards side by side
card_data = [
    ("Scope 1", "Direct emissions (own sources)", "LIVE", SUCCESS),
    ("Scope 2", "Indirect from purchased energy", "LIVE", SUCCESS),
    ("Scope 3", "Indirect from value chain", "COMING SOON", WARNING),
]
for i, (scope, desc, status, status_color) in enumerate(card_data):
    left = Inches(0.7 + i * 4.1)
    # Card background
    card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(1.8), Inches(3.7), Inches(3.5))
    card.fill.solid(); card.fill.fore_color.rgb = FIREFLY_50
    card.line.color.rgb = FIREFLY_100; card.line.width = Pt(1)
    card.adjustments[0] = 0.08
    # Scope label
    txB = s.shapes.add_textbox(left + Inches(0.3), Inches(2.1), Inches(3.1), Inches(0.5))
    tf = txB.text_frame; p = tf.paragraphs[0]
    r = p.add_run(); r.text = scope; r.font.size = Pt(20); r.font.color.rgb = FIREFLY_700; r.font.name = FONT; r.font.bold = True
    # Description
    txB2 = s.shapes.add_textbox(left + Inches(0.3), Inches(2.7), Inches(3.1), Inches(1.0))
    tf2 = txB2.text_frame; tf2.word_wrap = True; p2 = tf2.paragraphs[0]
    r2 = p2.add_run(); r2.text = desc; r2.font.size = Pt(13); r2.font.color.rgb = GRAY_700; r2.font.name = FONT
    # Status pill
    pill = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left + Inches(0.3), Inches(4.2), Inches(2.2), Inches(0.4))
    pill.fill.solid(); pill.fill.fore_color.rgb = status_color
    pill.line.fill.background()
    pill.adjustments[0] = 0.5
    tf_pill = pill.text_frame; tf_pill.margin_left = Pt(4); tf_pill.margin_right = Pt(4); tf_pill.margin_top = Pt(0); tf_pill.margin_bottom = Pt(0)
    p_pill = tf_pill.paragraphs[0]; p_pill.alignment = PP_ALIGN.CENTER
    r_pill = p_pill.add_run(); r_pill.text = status; r_pill.font.size = Pt(11); r_pill.font.color.rgb = WHITE; r_pill.font.name = FONT; r_pill.font.bold = True

add_note(s, "No specific dates \u2014 say \u2018it\u2019s in our active roadmap\u2019", top=Inches(6.0))

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 18 — When to Pitch Target Tracking
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide(prs)
add_title(s, "Conversation triggers \u2014 when to bring up TT", highlight_word="triggers")
add_section_label(s, "CLIENT PHRASES TO LISTEN FOR")
add_bullets(s, [
    "\u201cWe have published reduction goals but don\u2019t know if we\u2019re on track\u201d",
    "\u201cOur quarterly review is built in Excel and takes weeks\u201d",
    "\u201cThe board asks us to justify which projects to invest in\u201d",
    "\u201cWe need to report to CDP/SBTi/CSRD with auditable data\u201d",
    "\u201cWe want to see the impact of our projects on our goals\u201d",
    "\u201cWe have corporate goals but sites have no accountability\u201d",
    "Any mention of: \u2018environmental goals\u2019, \u2018reduction targets\u2019, \u2018net zero\u2019",
], top=Inches(1.7), font_size=Pt(13))

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 19 — Ideal Customer Profile
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide(prs)
add_title(s, "Who buys Target Tracking", highlight_word="Target Tracking")
add_section_label(s, "IDEAL CUSTOMER PROFILE")
add_bullets(s, [
    "Industry: Manufacturing, Food & Beverage, Consumer Goods, Mining",
    "Size: 5+ sites (multi-site value)",
    "Maturity: Already collecting environmental data (water and/or energy)",
    "Urgency: Published goals with deadlines (SBTi, CDP, annual report)",
    "Strong signal: Currently using Excel/Domo/Tableau to track goals",
], top=Inches(1.7))

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 20 — Common Objections (table)
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide(prs)
add_title(s, "The 5 objections you will hear \u2014 and what to say", highlight_word="what to say")
objections_data = [
    ["Objection", "Quick Response"],
    ["\u201cWe already do this in Excel\u201d", "\u201cHow long does your quarterly review take? Can you model 3 scenarios in 5 minutes?\u201d"],
    ["\u201cWe don\u2019t have energy/carbon data yet\u201d", "\u201cStart with water targets \u2014 add carbon when ready\u201d"],
    ["\u201cToo expensive for another module\u201d", "\u201cShow them the MACC curve \u2014 TT saves money by showing where NOT to invest\u201d"],
    ["\u201cOur sites are too different\u201d", "\u201cExactly: independent targets per site with consolidated view\u201d"],
    ["\u201cWe already use another carbon tool\u201d", "\u201cTT doesn\u2019t replace GHG inventory \u2014 it adds target tracking and scenario modeling\u201d"],
]
add_styled_table(s, objections_data, Inches(0.7), Inches(1.5), Inches(11.8), Inches(5.0))

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 21 — How to Start the Conversation
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide(prs)
add_title(s, "The question that opens the door", highlight_word="opens the door")

cards_data = [
    ("Existing clients", "You already have the data in the platform. Now you can turn it into targets with scenarios and automatic tracking. Can I show you in 15 minutes?"),
    ("New prospects", "How are you tracking progress against your water/carbon/energy goals?"),
    ("Upsell", "Did you know you can model investment scenarios and see which projects put you on track?"),
]
for i, (label, text) in enumerate(cards_data):
    left = Inches(0.7 + i * 4.1)
    card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(1.6), Inches(3.7), Inches(3.5))
    card.fill.solid(); card.fill.fore_color.rgb = FIREFLY_50
    card.line.color.rgb = FIREFLY_100; card.line.width = Pt(1)
    card.adjustments[0] = 0.08
    # Label pill
    pill = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left + Inches(0.3), Inches(1.9), Inches(2.5), Inches(0.35))
    pill.fill.solid(); pill.fill.fore_color.rgb = FIREFLY_500
    pill.line.fill.background(); pill.adjustments[0] = 0.5
    tf_pill = pill.text_frame; tf_pill.margin_left = Pt(4); tf_pill.margin_right = Pt(4); tf_pill.margin_top = Pt(0); tf_pill.margin_bottom = Pt(0)
    p_pill = tf_pill.paragraphs[0]; p_pill.alignment = PP_ALIGN.CENTER
    r_pill = p_pill.add_run(); r_pill.text = label.upper(); r_pill.font.size = Pt(10); r_pill.font.color.rgb = WHITE; r_pill.font.name = FONT; r_pill.font.bold = True
    # Text
    txB = s.shapes.add_textbox(left + Inches(0.3), Inches(2.5), Inches(3.1), Inches(2.3))
    tf = txB.text_frame; tf.word_wrap = True; p = tf.paragraphs[0]
    r = p.add_run(); r.text = f"\u201c{text}\u201d"
    r.font.size = Pt(12); r.font.color.rgb = GRAY_900; r.font.name = FONT

# Golden rule
golden = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(5.5), Inches(11.8), Inches(0.8))
golden.fill.solid(); golden.fill.fore_color.rgb = FIREFLY_100
golden.line.fill.background(); golden.adjustments[0] = 0.15
tf_g = golden.text_frame; tf_g.word_wrap = True
tf_g.margin_left = Pt(12); tf_g.margin_right = Pt(12); tf_g.margin_top = Pt(8); tf_g.margin_bottom = Pt(8)
p_g = tf_g.paragraphs[0]; p_g.alignment = PP_ALIGN.CENTER
r_g = p_g.add_run()
r_g.text = "Golden rule: Never start with \u2018we have a Target Tracking module\u2019. Always start with the client\u2019s problem."
r_g.font.size = Pt(13); r_g.font.color.rgb = FIREFLY_900; r_g.font.name = FONT; r_g.font.bold = True

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 22 — Cheat Sheet: Scope 1, 2, 3
# ═══════════════════════════════════════════════════════════════════════════
s = new_slide(prs)
add_title(s, "Quick reference: Scope 1, 2, and 3", highlight_word="Scope 1, 2, and 3")
scope_data = [
    ["Scope", "What it is", "Examples", "In Waterplan"],
    ["Scope 1", "Direct emissions (own sources)", "Natural gas, diesel, LPG", "Supported (Thermal variables)"],
    ["Scope 2", "Indirect from purchased energy", "Electricity, steam, RECs", "Supported (Electricity variables)"],
    ["Scope 3", "Indirect from value chain", "Logistics, raw materials", "Coming soon"],
]
add_styled_table(s, scope_data, Inches(0.7), Inches(1.5), Inches(11.8), Inches(3.0))

# Formula
txBox = s.shapes.add_textbox(Inches(0.7), Inches(5.0), Inches(11.8), Inches(0.8))
tf = txBox.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
# Label
r_label = p.add_run(); r_label.text = "Formula:  "
r_label.font.size = Pt(14); r_label.font.color.rgb = GRAY_700; r_label.font.name = FONT
# Formula text
r_formula = p.add_run(); r_formula.text = "Emissions = Activity \u00d7 Emission Factor"
r_formula.font.size = Pt(18); r_formula.font.color.rgb = FIREFLY_700; r_formula.font.name = FONT; r_formula.font.bold = True

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 23 — Close & Q&A
# ═══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
bg = s.background; fill = bg.fill; fill.solid(); fill.fore_color.rgb = WHITE
add_waterplan_logo(s)
add_title(s, "Next Steps", highlight_word="Next Steps")

add_bullets(s, [
    "This deck is yours \u2014 reuse slides 1-15 for client pitches",
    "Access the sales demo platform (ask for credentials if needed)",
    "Use the demo script to practice step by step",
    "Action: identify 2-3 clients in your pipeline where TT is relevant",
], top=Inches(1.5), font_size=Pt(15))

# Q&A block
qa_box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3.5), Inches(4.0), Inches(6.3), Inches(1.2))
qa_box.fill.solid(); qa_box.fill.fore_color.rgb = FIREFLY_50
qa_box.line.color.rgb = FIREFLY_100; qa_box.line.width = Pt(1)
qa_box.adjustments[0] = 0.15
tf_qa = qa_box.text_frame; tf_qa.word_wrap = True
tf_qa.margin_top = Pt(8); tf_qa.margin_bottom = Pt(8)
p_qa = tf_qa.paragraphs[0]; p_qa.alignment = PP_ALIGN.CENTER
r_qa = p_qa.add_run(); r_qa.text = "Questions?"
r_qa.font.size = Pt(28); r_qa.font.color.rgb = FIREFLY_700; r_qa.font.name = FONT; r_qa.font.bold = True

# Feedback note
add_note(s, "Anonymous feedback form coming after this session", top=Inches(5.8), left=Inches(3.5))

# ── Save ──────────────────────────────────────────────────────────────────
output_path = "/Users/ivokalaizicwp/Documents/waterplan/Q2-2026/Weekly planning/Week-May-5/04-Training-Target-Tracking/Target Tracking Training - May 2026.pptx"
os.makedirs(os.path.dirname(output_path), exist_ok=True)
prs.save(output_path)
print(f"Saved to: {output_path}")
print(f"Total slides: {len(prs.slides)}")
