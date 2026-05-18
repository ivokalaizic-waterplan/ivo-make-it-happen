"""Generate Q2 2026 Leadership Deck for Waterplan."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Brand Colors (Seastem/Firefly) ──
FIREFLY_500 = RGBColor(0x35, 0xA1, 0xC9)
FIREFLY_700 = RGBColor(0x20, 0x60, 0x78)
FIREFLY_900 = RGBColor(0x0F, 0x2D, 0x38)
FIREFLY_100 = RGBColor(0xD6, 0xEC, 0xF4)
FIREFLY_50  = RGBColor(0xEB, 0xF6, 0xFA)
GRAY_900    = RGBColor(0x2A, 0x2B, 0x2D)
GRAY_700    = RGBColor(0x55, 0x57, 0x5B)
GRAY_100    = RGBColor(0xDA, 0xDB, 0xDE)
GRAY_50     = RGBColor(0xEE, 0xEF, 0xF0)
WHITE       = RGBColor(0xFD, 0xFD, 0xFD)
SUCCESS     = RGBColor(0x02, 0x7A, 0x48)
WARNING     = RGBColor(0xB5, 0x47, 0x08)
ERROR       = RGBColor(0xB3, 0x23, 0x18)

FONT = "Inter"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]


# ── Helpers ──
def add_waterplan_logo(slide):
    logo_left = Inches(10.8)
    logo_top = Inches(0.3)
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, logo_left, logo_top, Inches(2), Inches(0.45)
    )
    shape.fill.background()
    shape.line.color.rgb = FIREFLY_500
    shape.line.width = Pt(1.5)
    shape.adjustments[0] = 0.5
    tf = shape.text_frame
    tf.margin_left = Pt(8); tf.margin_right = Pt(8)
    tf.margin_top = Pt(0); tf.margin_bottom = Pt(0)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "waterplan"
    run.font.size = Pt(13); run.font.color.rgb = FIREFLY_500
    run.font.name = FONT; run.font.bold = False


def add_title(slide, text, highlight_word=None):
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(10), Inches(0.8))
    tf = txBox.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    if highlight_word and highlight_word in text:
        parts = text.split(highlight_word, 1)
        for i, part in enumerate([parts[0], highlight_word, parts[1]]):
            if not part:
                continue
            run = p.add_run()
            run.text = part
            run.font.size = Pt(24); run.font.bold = True; run.font.name = FONT
            run.font.color.rgb = FIREFLY_500 if part == highlight_word else FIREFLY_700
    else:
        run = p.add_run()
        run.text = text
        run.font.size = Pt(24); run.font.bold = True; run.font.name = FONT
        run.font.color.rgb = FIREFLY_700


def add_subtitle(slide, text, top=Inches(1.0)):
    txBox = slide.shapes.add_textbox(Inches(0.5), top, Inches(12), Inches(0.5))
    tf = txBox.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(13); run.font.color.rgb = GRAY_700; run.font.name = FONT


def add_pill(slide, left, top, width, height, fill_color, text="",
             text_color=WHITE, font_size=11, bold=True):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid(); shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    shape.adjustments[0] = 0.25
    if text:
        tf = shape.text_frame; tf.word_wrap = False
        tf.margin_left = Pt(4); tf.margin_right = Pt(4)
        tf.margin_top = Pt(0); tf.margin_bottom = Pt(0)
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = text; run.font.size = Pt(font_size)
        run.font.color.rgb = text_color; run.font.name = FONT; run.font.bold = bold
    return shape


def add_card(slide, left, top, width, height, fill=FIREFLY_50, border=FIREFLY_100):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid(); shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = border; shape.line.width = Pt(1)
    shape.adjustments[0] = 0.08
    return shape


def add_text(slide, left, top, width, height, text, size=12, color=GRAY_900,
             bold=False, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    run = p.add_run()
    run.text = text; run.font.size = Pt(size); run.font.color.rgb = color
    run.font.name = FONT; run.font.bold = bold
    return tf


def add_section_label(slide, text, left, top):
    txBox = slide.shapes.add_textbox(left, top, Inches(3), Inches(0.3))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text.upper()
    run.font.size = Pt(11); run.font.color.rgb = FIREFLY_500
    run.font.name = FONT; run.font.bold = True


def add_footer(slide, text):
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(7.0), Inches(12), Inches(0.4))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(9); run.font.color.rgb = GRAY_700; run.font.name = FONT


# ═══════════════════════════════════════════════════════
# SLIDE 1 — Title
# ═══════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
# Background accent bar
bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(2.5), prs.slide_width, Inches(2.8))
bar.fill.solid(); bar.fill.fore_color.rgb = FIREFLY_900
bar.line.fill.background()

add_text(s, Inches(0.8), Inches(0.5), Inches(5), Inches(0.5),
         "waterplan", size=16, color=FIREFLY_500, bold=False)

add_text(s, Inches(0.8), Inches(2.8), Inches(10), Inches(0.8),
         "Q2 2026 Product Plan", size=36, color=WHITE, bold=True)
add_text(s, Inches(0.8), Inches(3.6), Inches(10), Inches(0.6),
         "Acciones del quarter para leadership", size=18, color=FIREFLY_100)
add_text(s, Inches(0.8), Inches(4.4), Inches(10), Inches(0.5),
         "Mayo - Junio - Julio 2026", size=14, color=FIREFLY_100)

add_footer(s, "Waterplan  |  Q2 2026  |  Confidential")


# ═══════════════════════════════════════════════════════
# SLIDE 2 — Guiding Principle + Rocks Overview
# ═══════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
add_waterplan_logo(s)
add_title(s, "Q2 North Star & Rocks", "North Star")

# North star card
add_card(s, Inches(0.5), Inches(1.3), Inches(12.3), Inches(1.0), fill=FIREFLY_50, border=FIREFLY_100)
add_text(s, Inches(0.8), Inches(1.35), Inches(11.5), Inches(0.35),
         "GUIDING PRINCIPLE", size=11, color=FIREFLY_500, bold=True)
add_text(s, Inches(0.8), Inches(1.7), Inches(11.5), Inches(0.5),
         "Que los sitios usen la plataforma. Every decision is filtered through: does this help site users adopt and trust the platform?",
         size=14, color=FIREFLY_900, bold=False)

# Three rock cards
rocks = [
    ("ROCK 1", "ABI — Smart Scenario\nPrioritization",
     "Discovery, design y MVP para que ABI priorice proyectos con constraints reales y AI.",
     "High"),
    ("ROCK 2", "Scope 3 Logistics\nDiscovery",
     "Full discovery: entender cómo calculan Scope 3, mapear metodología, definir diseño.",
     "Medium"),
    ("ROCK 3", "Sites Use\nTarget Tracking",
     "Permisos, homepage y side nav mobile para que cualquier cliente pueda hacer rollout a sitios.",
     "High"),
]

card_w = Inches(3.9)
card_h = Inches(3.6)
gap = Inches(0.3)
start_left = Inches(0.5)
card_top = Inches(2.7)

for i, (label, title, desc, complexity) in enumerate(rocks):
    left = start_left + i * (card_w + gap)
    add_card(s, left, card_top, card_w, card_h, fill=WHITE, border=GRAY_100)

    # Rock label pill
    add_pill(s, left + Inches(0.3), card_top + Inches(0.25), Inches(1.0), Inches(0.32),
             FIREFLY_500, label, WHITE, 10, True)

    # Complexity pill
    comp_color = ERROR if complexity == "High" else WARNING
    add_pill(s, left + Inches(1.4), card_top + Inches(0.25), Inches(1.0), Inches(0.32),
             comp_color, complexity, WHITE, 10, True)

    # Title
    add_text(s, left + Inches(0.3), card_top + Inches(0.75), Inches(3.3), Inches(1.0),
             title, size=16, color=FIREFLY_700, bold=True)

    # Description
    add_text(s, left + Inches(0.3), card_top + Inches(1.8), Inches(3.3), Inches(1.5),
             desc, size=12, color=GRAY_700)

add_footer(s, "Waterplan  |  Q2 2026  |  Rocks Overview")


# ═══════════════════════════════════════════════════════
# SLIDE 3 — Rock 1: ABI
# ═══════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
add_waterplan_logo(s)
add_title(s, "Rock 1: ABI — Smart Scenario Prioritization", "ABI")
add_subtitle(s, "Done = ABI can use the tool to prioritize projects with real constraints and get actionable recommendations.")

# Monthly milestones
months_data = [
    ("MAY", "Discovery", [
        "Map ABI's prioritization criteria, constraints, and decision process",
        "Deliverable: criteria document validated by ABI",
    ]),
    ("JUN", "Design + Architecture", [
        "Define AI agent approach, UX for constraint input and recommendations",
        "Deliverable: design approved by ABI, engineering ready to build",
    ]),
    ("JUL", "Build + Test", [
        "Engineering builds MVP, ABI tests with real data",
        "Deliverable: ABI runs a real prioritization through the tool",
    ]),
]

for i, (month, phase, items) in enumerate(months_data):
    top = Inches(1.8) + i * Inches(1.7)
    # Month pill
    add_pill(s, Inches(0.5), top, Inches(0.9), Inches(0.38), FIREFLY_700, month, WHITE, 12, True)
    # Phase
    add_text(s, Inches(1.6), top - Inches(0.02), Inches(3), Inches(0.4),
             phase, size=14, color=FIREFLY_900, bold=True)
    # Items
    for j, item in enumerate(items):
        add_text(s, Inches(1.6), top + Inches(0.38) + j * Inches(0.38), Inches(8), Inches(0.38),
                 "  " + item, size=12, color=GRAY_700)

# Risk card
risk_top = Inches(5.8)
add_card(s, Inches(0.5), risk_top, Inches(12.3), Inches(0.9), fill=RGBColor(0xFD, 0xF0, 0xE6), border=WARNING)
add_text(s, Inches(0.8), risk_top + Inches(0.05), Inches(1.5), Inches(0.35),
         "KEY RISKS", size=11, color=WARNING, bold=True)
add_text(s, Inches(0.8), risk_top + Inches(0.4), Inches(11.5), Inches(0.45),
         "Discovery doesn't converge (time-box 3-4 weeks)  |  AI architecture undefined (involve eng early)  |  ABI stakeholder availability",
         size=11, color=GRAY_700)

add_footer(s, "Waterplan  |  Q2 2026  |  Rock 1")


# ═══════════════════════════════════════════════════════
# SLIDE 4 — Rock 2: Scope 3 Logistics
# ═══════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
add_waterplan_logo(s)
add_title(s, "Rock 2: Scope 3 Logistics Discovery", "Scope 3")
add_subtitle(s, "Done = Methodology mapped + validated design. Engineering can plan implementation for Q3+.")

months_data = [
    ("MAY", "Current State", [
        "Understand current methodology: how do they calculate Scope 3 logistics today?",
        "Deliverable: methodology document with current-state mapping",
    ]),
    ("JUN", "Gap Analysis", [
        "Identify gaps and improvements: what can Waterplan do better?",
        "Deliverable: gap analysis + data requirements document",
    ]),
    ("JUL", "Design Validation", [
        "How this looks in the platform, what the UX is, what engineering needs",
        "Deliverable: design spec that engineering can estimate for Q3",
    ]),
]

for i, (month, phase, items) in enumerate(months_data):
    top = Inches(1.8) + i * Inches(1.7)
    add_pill(s, Inches(0.5), top, Inches(0.9), Inches(0.38), FIREFLY_700, month, WHITE, 12, True)
    add_text(s, Inches(1.6), top - Inches(0.02), Inches(3), Inches(0.4),
             phase, size=14, color=FIREFLY_900, bold=True)
    for j, item in enumerate(items):
        add_text(s, Inches(1.6), top + Inches(0.38) + j * Inches(0.38), Inches(8), Inches(0.38),
                 "  " + item, size=12, color=GRAY_700)

# Risk card
risk_top = Inches(5.8)
add_card(s, Inches(0.5), risk_top, Inches(12.3), Inches(0.9), fill=RGBColor(0xFD, 0xF0, 0xE6), border=WARNING)
add_text(s, Inches(0.8), risk_top + Inches(0.05), Inches(1.5), Inches(0.35),
         "KEY RISKS", size=11, color=WARNING, bold=True)
add_text(s, Inches(0.8), risk_top + Inches(0.4), Inches(11.5), Inches(0.45),
         "Methodology is complex/non-standard  |  Design dependency (involve design early)  |  Scope creep into implementation (Q2 = discovery only)",
         size=11, color=GRAY_700)

add_footer(s, "Waterplan  |  Q2 2026  |  Rock 2")


# ═══════════════════════════════════════════════════════
# SLIDE 5 — Rock 3: Sites Use Target Tracking
# ═══════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
add_waterplan_logo(s)
add_title(s, "Rock 3: Sites Use Target Tracking", "Target Tracking")
add_subtitle(s, "Done = Permission system configurable for both Coca-Cola and Colgate models. Homepage live. Side nav mobile.")

# Sub-initiatives table
rows = [
    ["Initiative", "Complexity", "Month"],
    ["Expand scenario limit (MAX_SCENARIOS)", "Very Low", "May (day 1)"],
    ["Homepage rollout for all customers", "Low-Med", "May"],
    ["Permission system design (full model)", "High", "May"],
    ["Role-based action permissions", "High", "June"],
    ["Configurable cross-site visibility", "High", "June"],
    ["Side nav mobile + flag cleanup", "Med-High", "June"],
    ["Lock permissions by role", "Low", "July"],
    ["Granular data visibility (hide financials)", "Medium", "July"],
    ["Breakdown by site respecting permissions", "Low", "July"],
]

table_shape = s.shapes.add_table(
    len(rows), 3, Inches(0.5), Inches(1.7), Inches(12.3), Inches(4.5)
)
table = table_shape.table
table.columns[0].width = Inches(6.5)
table.columns[1].width = Inches(2.5)
table.columns[2].width = Inches(3.3)

for row_idx, row in enumerate(rows):
    for col_idx, val in enumerate(row):
        cell = table.cell(row_idx, col_idx)
        cell.text = val
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        for paragraph in cell.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(11); run.font.name = FONT
        if row_idx == 0:
            cell.fill.solid(); cell.fill.fore_color.rgb = FIREFLY_700
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.color.rgb = WHITE; run.font.bold = True
        else:
            cell.fill.solid()
            cell.fill.fore_color.rgb = FIREFLY_50 if row_idx % 2 == 0 else WHITE
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.color.rgb = GRAY_900

add_footer(s, "Waterplan  |  Q2 2026  |  Rock 3")


# ═══════════════════════════════════════════════════════
# SLIDE 6 — Quarter Work (parallel initiatives)
# ═══════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
add_waterplan_logo(s)
add_title(s, "Quarter Work: Parallel Initiatives", "Parallel")
add_subtitle(s, "Run alongside the rocks. Important for the business but don't define the quarter's success.")

# Cards grid (2 columns x 3 rows)
initiatives = [
    ("Water/Carbon Separation", "High",
     "Separate Water and Carbon as independent entities in Target Tracking. Independent configs and scenario management per domain.",
     "Only affects Colgate today but needed for platform integrity."),
    ("Year-Aware Emission Factors", "Medium",
     "Year-aware emission + energy cost factors. Self-service management.",
     "Colgate will request 2025 data. Time-sensitive."),
    ("Unit Conversion", "Medium",
     "Unit conversion for water variables. Direct request from Tobias and Coca-Cola.",
     "Daily friction for active users."),
    ("Scenario Experience", "Medium",
     "Scenario-aware landing page + gap-to-target in absolute and financial terms.",
     "Trust impact. Budget justification enabler."),
    ("Projects Module", "Low-Med",
     "Multi-metric and multi-unit bulk upload. Continue Simplify initiative.",
     "Individual creation works. Gap is bulk upload only."),
    ("Data Access", "Medium",
     "Data download from UI for customers.",
     "Good value but needs scoping first."),
]

card_w = Inches(5.9)
card_h = Inches(1.55)
for i, (name, complexity, desc, why) in enumerate(initiatives):
    col = i % 2
    row = i // 2
    left = Inches(0.5) + col * (card_w + Inches(0.3))
    top = Inches(1.7) + row * (card_h + Inches(0.2))

    add_card(s, left, top, card_w, card_h, fill=WHITE, border=GRAY_100)

    # Name
    add_text(s, left + Inches(0.25), top + Inches(0.12), Inches(3.8), Inches(0.35),
             name, size=13, color=FIREFLY_700, bold=True)

    # Complexity pill
    comp_color = ERROR if complexity == "High" else (WARNING if "Med" in complexity else SUCCESS)
    add_pill(s, left + card_w - Inches(1.3), top + Inches(0.12), Inches(1.0), Inches(0.3),
             comp_color, complexity, WHITE, 9, True)

    # Description
    add_text(s, left + Inches(0.25), top + Inches(0.5), Inches(5.4), Inches(0.55),
             desc, size=11, color=GRAY_900)
    add_text(s, left + Inches(0.25), top + Inches(1.05), Inches(5.4), Inches(0.4),
             why, size=10, color=GRAY_700)

add_footer(s, "Waterplan  |  Q2 2026  |  Quarter Work")


# ═══════════════════════════════════════════════════════
# SLIDE 7 — Timeline
# ═══════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
add_waterplan_logo(s)
add_title(s, "Q2 Timeline: May - June - July", "Timeline")

months = ["MAY — Foundations", "JUNE — Build", "JULY — Deliver & Validate"]
month_colors = [FIREFLY_500, FIREFLY_700, FIREFLY_900]

col_w = Inches(3.9)
col_gap = Inches(0.3)

month_items = [
    [  # May
        ("Rock 1", "ABI discovery sessions"),
        ("Rock 2", "First discovery meetings"),
        ("Rock 3", "Expand scenario limit (day 1)"),
        ("Rock 3", "Homepage rollout (weeks 1-2)"),
        ("Rock 3", "Design permission system"),
        ("QW", "Start water/carbon separation"),
        ("QW", "Start year-aware emission factors"),
    ],
    [  # June
        ("Rock 1", "Design + architecture finalized"),
        ("Rock 2", "Gap analysis complete"),
        ("Rock 3", "Permissions MVP implemented"),
        ("Rock 3", "Side nav mobile delivered"),
        ("QW", "Water/carbon separation delivered"),
        ("QW", "Emission factors delivered"),
        ("QW", "Start unit conversion"),
    ],
    [  # July
        ("Rock 1", "MVP delivered, ABI tests"),
        ("Rock 2", "Design spec validated"),
        ("Rock 3", "Permissions extended + granular"),
        ("Rock 3", "2 clients configured (CC + CG)"),
        ("QW", "Unit conversion delivered"),
        ("QW", "Scenario-aware landing delivered"),
        ("QW", "Evaluate Tier 3 capacity"),
    ],
]

for i, (month_label, items) in enumerate(zip(months, month_items)):
    left = Inches(0.5) + i * (col_w + col_gap)

    # Month header pill
    add_pill(s, left, Inches(1.4), col_w, Inches(0.45), month_colors[i],
             month_label, WHITE, 13, True)

    # Items
    for j, (tag, desc) in enumerate(items):
        item_top = Inches(2.1) + j * Inches(0.65)

        # Tag pill
        tag_color = FIREFLY_500 if "Rock" in tag else GRAY_700
        add_pill(s, left + Inches(0.1), item_top, Inches(0.85), Inches(0.28),
                 tag_color, tag, WHITE, 8, True)

        # Desc
        add_text(s, left + Inches(1.05), item_top - Inches(0.02), Inches(2.75), Inches(0.35),
                 desc, size=11, color=GRAY_900)

add_footer(s, "Waterplan  |  Q2 2026  |  Timeline")


# ═══════════════════════════════════════════════════════
# SLIDE 8 — Risks & Mitigations
# ═══════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
add_waterplan_logo(s)
add_title(s, "Risk Summary & Mitigations", "Risk")

risks = [
    ("Water/Carbon separation touches ~145+ files",
     "High", "Start with backend model, migrate data early, frontend follows."),
    ("Permission system designed too narrow",
     "High", "Design full model upfront (Coca-Cola + Colgate + generic), implement incrementally."),
    ("ABI discovery doesn't converge",
     "Medium", "Time-box to 3-4 weeks. Define 'good enough' criteria early."),
    ("Colgate requests 2025 data before factors ready",
     "Medium", "Start emission factor work in May, don't wait for request."),
    ("Side nav mobile underestimated",
     "Medium", "Involve design early. Prototype on real devices."),
    ("Too many parallel initiatives",
     "High", "Rocks come first. Quarter work gets remaining capacity. If something slips, it's quarter work."),
]

for i, (risk, severity, mitigation) in enumerate(risks):
    top = Inches(1.4) + i * Inches(0.95)

    # Card background
    fill = RGBColor(0xFD, 0xF0, 0xE6) if severity == "High" else FIREFLY_50
    border = ERROR if severity == "High" else WARNING
    add_card(s, Inches(0.5), top, Inches(12.3), Inches(0.82), fill=fill, border=border)

    # Severity pill
    sev_color = ERROR if severity == "High" else WARNING
    add_pill(s, Inches(0.7), top + Inches(0.1), Inches(0.85), Inches(0.28),
             sev_color, severity, WHITE, 9, True)

    # Risk text
    add_text(s, Inches(1.7), top + Inches(0.05), Inches(5.5), Inches(0.35),
             risk, size=12, color=GRAY_900, bold=True)

    # Mitigation
    add_text(s, Inches(1.7), top + Inches(0.42), Inches(10.8), Inches(0.35),
             "Mitigation: " + mitigation, size=11, color=GRAY_700)

add_footer(s, "Waterplan  |  Q2 2026  |  Risk Summary")


# ── Save ──
output_dir = os.path.expanduser("~/Downloads")
output_path = os.path.join(output_dir, "waterplan-q2-2026-plan.pptx")
prs.save(output_path)
print(f"Saved to {output_path}")
