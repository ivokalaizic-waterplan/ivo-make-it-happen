"""
Q2 2026 Product Roadmap — Leadership Deck Generator
Uses Waterplan Firefly palette (internal deck)
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ─── Palette ───────────────────────────────────────────────────────────────────
FIREFLY_500 = RGBColor(0x35, 0xA1, 0xC9)
FIREFLY_700 = RGBColor(0x20, 0x60, 0x78)
FIREFLY_900 = RGBColor(0x0F, 0x2D, 0x38)
FIREFLY_100 = RGBColor(0xD6, 0xEC, 0xF4)
FIREFLY_50  = RGBColor(0xEB, 0xF6, 0xFA)
GRAY_900    = RGBColor(0x2A, 0x2B, 0x2D)
GRAY_700    = RGBColor(0x55, 0x57, 0x5B)
GRAY_500    = RGBColor(0x71, 0x73, 0x77)
GRAY_100    = RGBColor(0xDA, 0xDB, 0xDE)
GRAY_50     = RGBColor(0xEE, 0xEF, 0xF0)
WHITE       = RGBColor(0xFD, 0xFD, 0xFD)
SLIDE_BG    = RGBColor(0xF3, 0xF5, 0xF8)
SECTION_BG  = RGBColor(0xF0, 0xF2, 0xF5)
CARD_BG     = RGBColor(0xFF, 0xFF, 0xFF)

# Tier colors
TIER1_GREEN  = RGBColor(0x02, 0x7A, 0x48)
TIER1_GREEN_BG = RGBColor(0xE6, 0xF4, 0xED)
TIER2_YELLOW = RGBColor(0xB5, 0x47, 0x08)
TIER2_YELLOW_BG = RGBColor(0xFE, 0xF3, 0xE0)
DISCOVERY_BLUE = RGBColor(0x35, 0xA1, 0xC9)
DISCOVERY_BLUE_BG = RGBColor(0xEB, 0xF6, 0xFA)

# Complexity colors
COMPLEXITY_HIGH = RGBColor(0xB3, 0x23, 0x18)
COMPLEXITY_MED = RGBColor(0xB5, 0x47, 0x08)
COMPLEXITY_LOW = RGBColor(0x02, 0x7A, 0x48)

# Risk/warning
RED_700 = RGBColor(0xB3, 0x23, 0x18)
RED_50  = RGBColor(0xFE, 0xF3, 0xF2)

SKILL_DIR = os.path.expanduser("~/.claude/skills/waterplan-pptx")
ASSETS = os.path.join(SKILL_DIR, "assets")

# ─── Presentation Setup ────────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
blank_layout = prs.slide_layouts[6]


# ─── Helper Functions ──────────────────────────────────────────────────────────

def set_slide_bg(slide, color=SLIDE_BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_accent_bars(slide):
    """Top and bottom accent bars."""
    # Top bar
    top = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Pt(5)
    )
    top.fill.solid()
    top.fill.fore_color.rgb = FIREFLY_500
    top.line.fill.background()
    # Bottom bar
    bot = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(7.5) - Pt(5), Inches(13.333), Pt(5)
    )
    bot.fill.solid()
    bot.fill.fore_color.rgb = FIREFLY_500
    bot.line.fill.background()


def add_logo(slide):
    logo_path = os.path.join(ASSETS, "waterplan-logo-dark.png")
    if os.path.exists(logo_path):
        slide.shapes.add_picture(
            logo_path, Inches(11.3), Inches(0.35), Inches(1.6), Inches(0.32)
        )


def add_slide_title(slide, text, subtitle=None):
    txBox = slide.shapes.add_textbox(Inches(0.7), Inches(0.4), Inches(9), Inches(0.7))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(28)
    run.font.color.rgb = FIREFLY_900
    run.font.name = "Inter"
    run.font.bold = True

    if subtitle:
        p2 = tf.add_paragraph()
        run2 = p2.add_run()
        run2.text = subtitle
        run2.font.size = Pt(14)
        run2.font.color.rgb = GRAY_700
        run2.font.name = "Inter"


def add_card(slide, left, top, width, height, fill=CARD_BG):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = GRAY_100
    shape.line.width = Pt(0.75)
    shape.shadow.inherit = False
    return shape


def add_pill(slide, left, top, width, height, fill_color, text="",
             text_color=WHITE, font_size=11, bold=True):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    if text:
        tf = shape.text_frame
        tf.word_wrap = False
        tf.margin_left = Pt(6)
        tf.margin_right = Pt(6)
        tf.margin_top = Pt(2)
        tf.margin_bottom = Pt(2)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = text
        run.font.size = Pt(font_size)
        run.font.color.rgb = text_color
        run.font.name = "Inter"
        run.font.bold = bold
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=12,
                 color=GRAY_900, bold=False, alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(8)
    tf.margin_right = Pt(8)
    tf.margin_top = Pt(4)
    tf.margin_bottom = Pt(4)
    p = tf.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.color.rgb = color
    run.font.name = "Inter"
    run.font.bold = bold
    return txBox


def setup_slide(title, subtitle=None):
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide)
    add_accent_bars(slide)
    add_logo(slide)
    add_slide_title(slide, title, subtitle)
    return slide


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 1: Title
# ═══════════════════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, FIREFLY_900)

# Accent bars in lighter shade
top = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Pt(6)
)
top.fill.solid()
top.fill.fore_color.rgb = FIREFLY_500
top.line.fill.background()
bot = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(0), Inches(7.5) - Pt(6), Inches(13.333), Pt(6)
)
bot.fill.solid()
bot.fill.fore_color.rgb = FIREFLY_500
bot.line.fill.background()

# Logo (white version)
logo_path = os.path.join(ASSETS, "waterplan-logo-white.png")
if os.path.exists(logo_path):
    slide.shapes.add_picture(logo_path, Inches(0.7), Inches(0.5), Inches(2.0), Inches(0.4))

# Title
txBox = slide.shapes.add_textbox(Inches(0.7), Inches(2.5), Inches(11), Inches(1.5))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
run = p.add_run()
run.text = "Q2 2026 Product Roadmap"
run.font.size = Pt(44)
run.font.color.rgb = WHITE
run.font.name = "Inter"
run.font.bold = True

# Subtitle
p2 = tf.add_paragraph()
p2.space_before = Pt(12)
run2 = p2.add_run()
run2.text = "Leadership Review  |  May 2026"
run2.font.size = Pt(20)
run2.font.color.rgb = FIREFLY_100
run2.font.name = "Inter"

# Decorative line
line = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(0.7), Inches(4.3), Inches(3), Pt(4)
)
line.fill.solid()
line.fill.fore_color.rgb = FIREFLY_500
line.line.fill.background()

# Footer info
txBox = slide.shapes.add_textbox(Inches(0.7), Inches(5.8), Inches(6), Inches(0.8))
tf = txBox.text_frame
p = tf.paragraphs[0]
run = p.add_run()
run.text = "Product & Engineering  |  Ivo Kalaizic"
run.font.size = Pt(14)
run.font.color.rgb = FIREFLY_100
run.font.name = "Inter"


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 2: Personal Rocks
# ═══════════════════════════════════════════════════════════════════════════════

slide = setup_slide("Personal Rocks Q2 2026", "6 key outcomes I'm accountable for this quarter")

rocks = [
    ("R1", "ABI — Smart Scenario Prioritization",
     "Deliver MVP for AB InBev project prioritization",
     "Discovery: May  |  Design: June  |  MVP: Aug"),
    ("R2", "Colgate Site Rollout Success",
     "Site users actively using platform post-training",
     "Training: May 18-25  |  Blockers resolved: June"),
    ("R3", "Scope 3 Logistics Discovery",
     "Methodology mapped + validated design for Q3 eng",
     "Discovery: May-June  |  Design spec: July"),
    ("R4", "CCEP Site Adoption",
     "Gap analysis + pain point identified + action plan",
     "Gap analysis: May-June  |  Action plan: July"),
    ("R5", "Coca-Cola Platform Expansion",
     "Yearly Targets Calculator + True Cost of Water discovery",
     "Discovery: May-July  |  Design validated: End Q2"),
    ("R6", "Weekly Client Interviews",
     "1+ interview/week, insights documented & fed back",
     "Weekly, starting May — ongoing"),
]

# Section band
band = slide.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.4), Inches(12.333), Inches(5.7)
)
band.fill.solid()
band.fill.fore_color.rgb = SECTION_BG
band.line.fill.background()

for i, (num, title, objective, timeline) in enumerate(rocks):
    row = i // 2
    col = i % 2
    left = Inches(0.8) + col * Inches(6.2)
    top = Inches(1.6) + row * Inches(1.85)

    # Card
    card = add_card(slide, left, top, Inches(5.8), Inches(1.7))

    # Rock number pill
    add_pill(slide, left + Inches(0.2), top + Inches(0.2),
             Inches(0.5), Inches(0.35), FIREFLY_500, num, WHITE, 11, True)

    # Title
    add_text_box(slide, left + Inches(0.85), top + Inches(0.12),
                 Inches(4.7), Inches(0.4), title, 13, FIREFLY_900, True)

    # Objective
    add_text_box(slide, left + Inches(0.2), top + Inches(0.55),
                 Inches(5.4), Inches(0.5), objective, 11, GRAY_700)

    # Timeline
    add_text_box(slide, left + Inches(0.2), top + Inches(1.1),
                 Inches(5.4), Inches(0.4), timeline, 10, FIREFLY_700, True)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 3: Initiative Priority Matrix
# ═══════════════════════════════════════════════════════════════════════════════

slide = setup_slide("Initiative Priority Matrix", "14 initiatives ordered by priority — color-coded by tier")

initiatives = [
    ("P1", "ABI — Smart Scenario Prioritization", "High", "May → Sep", "tier1"),
    ("P2", "Permissions — Adapt TT to Roles", "Medium", "May", "tier1"),
    ("P3", "Water/Carbon Scenario Separation", "High", "May → Jul", "tier1"),
    ("P4", "Year-Aware Emission Factors", "Medium", "May-Jun (cond.)", "tier1"),
    ("P5", "Homepage Rollout (All Customers)", "Low-Med", "May", "tier1"),
    ("P6", "Target Progress & Scenario UX", "Medium", "May → Jul", "tier1"),
    ("P7", "Unit Conversion for Water", "Medium", "Jun-Jul", "tier1"),
    ("P8", "Cross-Site Visibility", "High", "May → Jul", "tier1"),
    ("P9", "Side Navigation — Mobile + Cleanup", "Med-High", "Jun-Jul", "tier2"),
    ("P10", "Simplify TT + Projects Nav", "TBD", "May → TBD", "tier2"),
    ("P11", "Scaling Carbon Financial Features", "Medium", "Conditional", "tier2"),
    ("D1", "Scope 3 Logistics — Discovery", "Low", "May-Jul", "discovery"),
    ("D2", "CCEP Yearly Targets — Discovery", "Low", "Q2", "discovery"),
    ("D3", "True Cost of Water — Discovery", "Low", "Q2", "discovery"),
]

# Legend pills
add_pill(slide, Inches(0.7), Inches(1.25), Inches(1.6), Inches(0.3),
         TIER1_GREEN, "TIER 1 — Must Do", WHITE, 9)
add_pill(slide, Inches(2.5), Inches(1.25), Inches(1.8), Inches(0.3),
         TIER2_YELLOW, "TIER 2 — Should Do", WHITE, 9)
add_pill(slide, Inches(4.5), Inches(1.25), Inches(1.5), Inches(0.3),
         DISCOVERY_BLUE, "DISCOVERY", WHITE, 9)

# Table header
header_top = Inches(1.75)
col_lefts = [Inches(0.7), Inches(1.3), Inches(7.0), Inches(9.2)]
col_widths = [Inches(0.6), Inches(5.5), Inches(2.0), Inches(2.5)]
headers = ["#", "Initiative", "Complexity", "Timeline"]

# Header row
header_band = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(0.5), header_top, Inches(12.333), Inches(0.4)
)
header_band.fill.solid()
header_band.fill.fore_color.rgb = FIREFLY_700
header_band.line.fill.background()

for idx, h in enumerate(headers):
    add_text_box(slide, col_lefts[idx], header_top, col_widths[idx], Inches(0.4),
                 h, 11, WHITE, True)

# Rows
row_h = Inches(0.35)
for i, (pri, name, complexity, timeline, tier) in enumerate(initiatives):
    top = header_top + Inches(0.45) + i * row_h

    # Alternating background
    if i % 2 == 0:
        row_bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0.5), top, Inches(12.333), row_h
        )
        row_bg.fill.solid()
        if tier == "tier1":
            row_bg.fill.fore_color.rgb = TIER1_GREEN_BG
        elif tier == "tier2":
            row_bg.fill.fore_color.rgb = TIER2_YELLOW_BG
        else:
            row_bg.fill.fore_color.rgb = DISCOVERY_BLUE_BG
        row_bg.line.fill.background()

    # Priority label
    tier_color = TIER1_GREEN if tier == "tier1" else (TIER2_YELLOW if tier == "tier2" else DISCOVERY_BLUE)
    add_text_box(slide, col_lefts[0], top, col_widths[0], row_h, pri, 10, tier_color, True)
    add_text_box(slide, col_lefts[1], top, col_widths[1], row_h, name, 10, GRAY_900, False)
    # Complexity color
    comp_color = COMPLEXITY_HIGH if "High" in complexity else (COMPLEXITY_MED if "Med" in complexity else COMPLEXITY_LOW)
    add_text_box(slide, col_lefts[2], top, col_widths[2], row_h, complexity, 10, comp_color, True)
    add_text_box(slide, col_lefts[3], top, col_widths[3], row_h, timeline, 10, GRAY_700, False)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 4: Tier 1 Deep Dive (P1-P4)
# ═══════════════════════════════════════════════════════════════════════════════

slide = setup_slide("Tier 1 Initiatives — Deep Dive (1/2)", "Must-do: committed deals & foundational work")

tier1_a = [
    ("P1", "ABI — Smart Scenario Prioritization", "HIGH", COMPLEXITY_HIGH,
     "Committed sale. 4-phase tool:\nUpload → Validate → Prioritize → KPIs",
     "Pilot live Aug-Sep (SAS zone: Brazil + LATAM)",
     "Discovery: May  |  Design+Arch: June  |  Build: July  |  Pilot: Aug-Sep"),
    ("P2", "Permissions — Adapt TT to Roles", "MEDIUM", COMPLEXITY_MED,
     "Map TT actions to existing Waterplan role system.\nCritical for enterprise rollout.",
     "Must ship before May 18 Colgate training",
     "May (hard deadline: May 18)"),
    ("P3", "Water/Carbon Scenario Separation", "HIGH", COMPLEXITY_HIGH,
     "Make Water & Carbon truly independent.\nEach domain gets own scenarios + metrics.",
     "Foundational for site adoption — largest arch change of Q2",
     "Backend: May-Jun  |  Frontend: Jun-Jul"),
    ("P4", "Year-Aware Emission Factors", "MEDIUM", COMPLEXITY_MED,
     "Add year dimension to emission factors.\nLookup by year with fallback logic.",
     "CONDITIONAL — only if Colgate sends 2025 data",
     "If needed: May-Jun"),
]

band = slide.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.3), Inches(12.333), Inches(5.85)
)
band.fill.solid()
band.fill.fore_color.rgb = SECTION_BG
band.line.fill.background()

for i, (pri, title, complexity, comp_color, desc, impact, timeline) in enumerate(tier1_a):
    row = i // 2
    col = i % 2
    left = Inches(0.7) + col * Inches(6.2)
    top = Inches(1.5) + row * Inches(2.9)

    card = add_card(slide, left, top, Inches(5.9), Inches(2.7))

    # Priority + Complexity pills
    add_pill(slide, left + Inches(0.2), top + Inches(0.2),
             Inches(0.5), Inches(0.3), FIREFLY_500, pri, WHITE, 10)
    add_pill(slide, left + Inches(0.8), top + Inches(0.2),
             Inches(0.9), Inches(0.3), comp_color, complexity, WHITE, 9)

    # Title
    add_text_box(slide, left + Inches(0.2), top + Inches(0.6),
                 Inches(5.5), Inches(0.35), title, 13, FIREFLY_900, True)

    # Description
    add_text_box(slide, left + Inches(0.2), top + Inches(0.95),
                 Inches(5.5), Inches(0.7), desc, 10, GRAY_700)

    # Impact
    add_text_box(slide, left + Inches(0.2), top + Inches(1.7),
                 Inches(5.5), Inches(0.45), impact, 10, FIREFLY_700, True)

    # Timeline bar
    tl_shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left + Inches(0.2), top + Inches(2.2),
        Inches(5.5), Inches(0.35)
    )
    tl_shape.fill.solid()
    tl_shape.fill.fore_color.rgb = FIREFLY_50
    tl_shape.line.fill.background()
    add_text_box(slide, left + Inches(0.3), top + Inches(2.2),
                 Inches(5.3), Inches(0.35), timeline, 9, FIREFLY_700, False)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 5: Tier 1 Deep Dive (P5-P8)
# ═══════════════════════════════════════════════════════════════════════════════

slide = setup_slide("Tier 1 Initiatives — Deep Dive (2/2)", "Quick wins + customer-driven requests")

tier1_b = [
    ("P5", "Homepage Rollout (All Customers)", "LOW-MED", COMPLEXITY_LOW,
     "Configurable homepage per company.\nDesign ready, component exists.",
     "Quick win — improves first impression & adoption",
     "May (quick win)"),
    ("P6", "Target Progress & Scenario Experience", "MEDIUM", COMPLEXITY_MED,
     "3 sub-deliverables:\n1. Settings panel  2. Gap-to-target  3. Scenario landing",
     "Directly affects how users perceive progress",
     "Settings: May  |  Gap: Jun  |  Landing: Jul"),
    ("P7", "Unit Conversion for Water", "MEDIUM", COMPLEXITY_MED,
     "Allow m³ display for Coca-Cola users.\nConversion table + selection UI + display layer.",
     "Direct & repeated customer request — eliminates daily friction",
     "Jun-Jul"),
    ("P8", "Cross-Site Visibility", "HIGH", COMPLEXITY_HIGH,
     "Define data visibility rules across sites.\nPending validation — may simplify scope.",
     "Pending validation with Tobías (client)",
     "Validation: May  |  If confirmed: Jun-Jul"),
]

band = slide.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.3), Inches(12.333), Inches(5.85)
)
band.fill.solid()
band.fill.fore_color.rgb = SECTION_BG
band.line.fill.background()

for i, (pri, title, complexity, comp_color, desc, impact, timeline) in enumerate(tier1_b):
    row = i // 2
    col = i % 2
    left = Inches(0.7) + col * Inches(6.2)
    top = Inches(1.5) + row * Inches(2.9)

    card = add_card(slide, left, top, Inches(5.9), Inches(2.7))

    add_pill(slide, left + Inches(0.2), top + Inches(0.2),
             Inches(0.5), Inches(0.3), FIREFLY_500, pri, WHITE, 10)
    add_pill(slide, left + Inches(0.8), top + Inches(0.2),
             Inches(1.0), Inches(0.3), comp_color, complexity, WHITE, 9)

    add_text_box(slide, left + Inches(0.2), top + Inches(0.6),
                 Inches(5.5), Inches(0.35), title, 13, FIREFLY_900, True)
    add_text_box(slide, left + Inches(0.2), top + Inches(0.95),
                 Inches(5.5), Inches(0.7), desc, 10, GRAY_700)
    add_text_box(slide, left + Inches(0.2), top + Inches(1.7),
                 Inches(5.5), Inches(0.45), impact, 10, FIREFLY_700, True)

    tl_shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left + Inches(0.2), top + Inches(2.2),
        Inches(5.5), Inches(0.35)
    )
    tl_shape.fill.solid()
    tl_shape.fill.fore_color.rgb = FIREFLY_50
    tl_shape.line.fill.background()
    add_text_box(slide, left + Inches(0.3), top + Inches(2.2),
                 Inches(5.3), Inches(0.35), timeline, 9, FIREFLY_700, False)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 6: Discovery Initiatives
# ═══════════════════════════════════════════════════════════════════════════════

slide = setup_slide("Discovery Initiatives", "No engineering in Q2 — research & design to prepare Q3 execution")

discoveries = [
    ("D1", "Scope 3 Logistics", "Ivo + Pipe",
     "Understand how clients calculate Scope 3 logistics emissions.\nMap methodology (GLEC, transport modes).",
     "Output: Methodology doc + data model design + validated spec",
     "Discovery: May-Jun  |  Design: Jul  |  Eng: Q3+"),
    ("D2", "CCEP Yearly Targets Calculator", "Ivo + Cami",
     "Coca-Cola uses ~90 Excel files for annual targets (15+ years old).\nValidate if Waterplan can absorb this process.",
     "Output: Methodology mapped + generic product opportunity validated",
     "Discovery: Q2  |  Design: End Q2  |  Eng: Q3+"),
    ("D3", "True Cost of Water", "Ivo + Cami",
     "Sites calculate real cost of water in separate spreadsheets.\nValidate if calculation is standard across sites.",
     "Output: Process documented + standardization assessment",
     "Discovery: Q2  |  Eng: If validated, Q3+"),
]

band = slide.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.4), Inches(12.333), Inches(5.6)
)
band.fill.solid()
band.fill.fore_color.rgb = SECTION_BG
band.line.fill.background()

for i, (code, title, owner, desc, output, timeline) in enumerate(discoveries):
    left = Inches(0.8) + i * Inches(4.1)
    top = Inches(1.7)

    card = add_card(slide, left, top, Inches(3.8), Inches(5.0))

    # Code pill
    add_pill(slide, left + Inches(0.2), top + Inches(0.2),
             Inches(0.5), Inches(0.3), DISCOVERY_BLUE, code, WHITE, 10)
    # Owner pill
    add_pill(slide, left + Inches(0.8), top + Inches(0.2),
             Inches(1.4), Inches(0.3), GRAY_100, owner, GRAY_700, 9, False)

    # Title
    add_text_box(slide, left + Inches(0.2), top + Inches(0.65),
                 Inches(3.4), Inches(0.4), title, 14, FIREFLY_900, True)

    # Description
    add_text_box(slide, left + Inches(0.2), top + Inches(1.1),
                 Inches(3.4), Inches(1.2), desc, 10, GRAY_700)

    # Output section
    add_text_box(slide, left + Inches(0.2), top + Inches(2.6),
                 Inches(3.4), Inches(0.25), "OUTPUT", 9, FIREFLY_500, True)
    add_text_box(slide, left + Inches(0.2), top + Inches(2.85),
                 Inches(3.4), Inches(0.8), output, 10, GRAY_900)

    # Timeline
    tl = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left + Inches(0.2), top + Inches(4.2),
        Inches(3.4), Inches(0.5)
    )
    tl.fill.solid()
    tl.fill.fore_color.rgb = DISCOVERY_BLUE_BG
    tl.line.fill.background()
    add_text_box(slide, left + Inches(0.3), top + Inches(4.25),
                 Inches(3.2), Inches(0.4), timeline, 9, FIREFLY_700, False)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 7: Timeline / Roadmap
# ═══════════════════════════════════════════════════════════════════════════════

slide = setup_slide("Q2 2026 Roadmap Timeline", "Month-by-month delivery view")

months = ["MAY", "JUNE", "JULY", "AUG-SEP"]
month_left = Inches(3.0)
month_width = Inches(2.4)

# Month headers
for i, m in enumerate(months):
    left = month_left + i * month_width
    add_pill(slide, left, Inches(1.5), Inches(2.1), Inches(0.35),
             FIREFLY_500, m, WHITE, 11)

# Workstream rows
workstreams = [
    ("ABI Prioritization", [
        (0, 1, "Discovery + Mock", FIREFLY_100),
        (1, 2, "Design + Arch + Build", FIREFLY_500),
        (3, 4, "Pilot Live", TIER1_GREEN),
    ]),
    ("Permissions (Actions)", [
        (0, 1, "Ship guards (before May 18)", TIER1_GREEN),
    ]),
    ("Water/Carbon Separation", [
        (0, 2, "Backend domain + Migration", FIREFLY_500),
        (1, 3, "Frontend + Settings", FIREFLY_100),
    ]),
    ("Homepage + Settings", [
        (0, 1, "Homepage + Settings panel", TIER1_GREEN),
    ]),
    ("Target Progress UX", [
        (1, 2, "Gap-to-target", FIREFLY_500),
        (2, 3, "Scenario landing", FIREFLY_100),
    ]),
    ("Unit Conversion", [
        (1, 3, "Table + UI + Display", FIREFLY_500),
    ]),
    ("Cross-Site Visibility", [
        (0, 1, "Validate w/ client", TIER2_YELLOW_BG),
        (1, 3, "Design + Implement (if confirmed)", FIREFLY_100),
    ]),
    ("Side Nav + Simplify", [
        (1, 3, "Mobile design + Cleanup", TIER2_YELLOW_BG),
    ]),
    ("Discovery (Scope 3 / CCEP / TCW)", [
        (0, 3, "Research + Methodology + Design", DISCOVERY_BLUE_BG),
    ]),
]

row_top = Inches(2.1)
row_h = Inches(0.55)

for i, (label, bars) in enumerate(workstreams):
    y = row_top + i * row_h

    # Label
    add_text_box(slide, Inches(0.3), y, Inches(2.7), row_h, label, 9, GRAY_900, False)

    # Separator line
    sep = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(3.0), y + row_h - Pt(1), Inches(9.8), Pt(1)
    )
    sep.fill.solid()
    sep.fill.fore_color.rgb = GRAY_100
    sep.line.fill.background()

    # Bars
    for start, end, bar_text, color in bars:
        bar_left = month_left + start * month_width + Inches(0.05)
        bar_width = (end - start) * month_width - Inches(0.1)
        bar = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, bar_left, y + Inches(0.08),
            bar_width, Inches(0.35)
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = color
        bar.line.fill.background()

        # Bar text
        tf = bar.text_frame
        tf.word_wrap = False
        tf.margin_left = Pt(6)
        tf.margin_top = Pt(1)
        tf.margin_bottom = Pt(1)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = bar_text
        run.font.size = Pt(8)
        run.font.color.rgb = GRAY_900
        run.font.name = "Inter"
        run.font.bold = False


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 8: Key Risks & Dependencies
# ═══════════════════════════════════════════════════════════════════════════════

slide = setup_slide("Key Risks & Dependencies")

risks = [
    ("HARD DEADLINE", "P2 Permissions must ship before May 18",
     "Colgate training starts — site users need role-based access on day one.",
     RED_700),
    ("CONDITIONAL", "P4 Year-Aware Factors depends on Colgate data",
     "If Colgate doesn't send 2025 emissions data this quarter, this defers to Q3.",
     TIER2_YELLOW),
    ("PENDING VALIDATION", "P8 Cross-Site Visibility awaiting client input",
     "Scope may shrink significantly if alternative approach (active projects view) is sufficient.",
     TIER2_YELLOW),
    ("TIME-BOX RISK", "ABI discovery must converge in 3-4 weeks",
     "Risk of scope creep. Multiple ABI stakeholders need to align on prioritization criteria.",
     COMPLEXITY_MED),
    ("CONDITIONAL", "P11 Carbon Financial Features depends on pipeline",
     "Only invest if new carbon customers materialize. Otherwise defer entirely.",
     GRAY_700),
]

band = slide.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.3), Inches(12.333), Inches(5.8)
)
band.fill.solid()
band.fill.fore_color.rgb = SECTION_BG
band.line.fill.background()

for i, (badge, title, desc, badge_color) in enumerate(risks):
    top = Inches(1.5) + i * Inches(1.1)
    card = add_card(slide, Inches(0.8), top, Inches(11.7), Inches(0.95))

    # Badge
    add_pill(slide, Inches(1.0), top + Inches(0.15),
             Inches(1.8), Inches(0.3), badge_color, badge, WHITE, 9)

    # Title
    add_text_box(slide, Inches(3.0), top + Inches(0.08),
                 Inches(9.0), Inches(0.35), title, 12, GRAY_900, True)

    # Description
    add_text_box(slide, Inches(3.0), top + Inches(0.45),
                 Inches(9.0), Inches(0.4), desc, 10, GRAY_700)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 9: Team Allocation
# ═══════════════════════════════════════════════════════════════════════════════

slide = setup_slide("Team Allocation", "Who is focused on what this quarter")

team = [
    ("Ivo", "Product & Discovery Lead", FIREFLY_700,
     ["ABI discovery (with Gustavo)", "Scope 3 discovery (with Pipe)",
      "CCEP + True Cost (with Cami)", "Permissions product spec",
      "Weekly client interviews"]),
    ("Ema", "Engineering Lead", FIREFLY_500,
     ["ABI phases 1-4 (arch + build)", "Water/Carbon separation (backend)",
      "Year-aware factors", "Unit conversion (backend)",
      "Cross-site visibility (if confirmed)"]),
    ("Mica", "Frontend Engineering", RGBColor(0x4B, 0x9F, 0x7F),
     ["Homepage rollout (quick win)", "Settings panel + scenario UX",
      "Side nav mobile + cleanup", "Unit conversion UI",
      "ABI frontend (with Ema)"]),
    ("Maria", "Design", RGBColor(0x8B, 0x5C, 0xF6),
     ["ABI UX design", "Mobile nav design",
      "Scope 3 design spec (Jul)", "Projects simplification"]),
    ("Sofi", "QA", GRAY_700,
     ["Permissions testing", "Role combinations validation"]),
    ("Pipe", "Discovery", GRAY_700,
     ["Scope 3 methodology research"]),
    ("Cami", "Discovery", GRAY_700,
     ["CCEP yearly targets", "True Cost of Water"]),
]

# Main team (top row - 4 cards)
for i in range(4):
    name, role, color, tasks = team[i]
    left = Inches(0.6) + i * Inches(3.15)
    top = Inches(1.4)

    card = add_card(slide, left, top, Inches(3.0), Inches(3.6))

    # Name pill
    add_pill(slide, left + Inches(0.15), top + Inches(0.15),
             Inches(1.2), Inches(0.35), color, name, WHITE, 11)

    # Role
    add_text_box(slide, left + Inches(0.15), top + Inches(0.55),
                 Inches(2.7), Inches(0.3), role, 9, GRAY_500)

    # Tasks
    task_text = "\n".join(f"• {t}" for t in tasks)
    txBox = slide.shapes.add_textbox(left + Inches(0.15), top + Inches(0.9),
                                      Inches(2.7), Inches(2.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(4)
    for j, task in enumerate(tasks):
        if j == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(3)
        run = p.add_run()
        run.text = f"• {task}"
        run.font.size = Pt(9)
        run.font.color.rgb = GRAY_700
        run.font.name = "Inter"

# Support team (bottom row - 3 cards)
for i in range(4, 7):
    name, role, color, tasks = team[i]
    col = i - 4
    left = Inches(0.6) + col * Inches(4.2)
    top = Inches(5.2)

    card = add_card(slide, left, top, Inches(3.9), Inches(1.8))

    add_pill(slide, left + Inches(0.15), top + Inches(0.15),
             Inches(1.0), Inches(0.3), color, name, WHITE, 10)

    add_text_box(slide, left + Inches(1.3), top + Inches(0.12),
                 Inches(2.4), Inches(0.3), role, 9, GRAY_500)

    task_text = "  |  ".join(tasks)
    add_text_box(slide, left + Inches(0.15), top + Inches(0.55),
                 Inches(3.6), Inches(1.1), "\n".join(f"• {t}" for t in tasks), 9, GRAY_700)


# ═══════════════════════════════════════════════════════════════════════════════
# Save
# ═══════════════════════════════════════════════════════════════════════════════

output_path = "/Users/ivokalaizicwp/Documents/waterplan/Q2-2026/Q2 - Roadmap/Q2 planning/outputs/Q2_2026_Roadmap_Leadership_Deck.pptx"
prs.save(output_path)
print(f"Saved: {output_path}")
