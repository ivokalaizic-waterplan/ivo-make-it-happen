#!/usr/bin/env python3
"""Generate E-Leads Workshop Executive Summary deck for Tobias."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ─── Palette (Firefly / default) ───────────────────────────────────
FIREFLY_500 = RGBColor(0x35, 0xA1, 0xC9)
FIREFLY_700 = RGBColor(0x20, 0x60, 0x78)
FIREFLY_900 = RGBColor(0x0F, 0x2D, 0x38)
FIREFLY_100 = RGBColor(0xD6, 0xEC, 0xF4)
FIREFLY_50  = RGBColor(0xEB, 0xF6, 0xFA)
GRAY_900    = RGBColor(0x2A, 0x2B, 0x2D)
GRAY_800    = RGBColor(0x3F, 0x41, 0x44)
GRAY_700    = RGBColor(0x55, 0x57, 0x5B)
GRAY_500    = RGBColor(0x80, 0x82, 0x86)
GRAY_100    = RGBColor(0xDA, 0xDB, 0xDE)
GRAY_50     = RGBColor(0xEE, 0xEF, 0xF0)
WHITE       = RGBColor(0xFD, 0xFD, 0xFD)
PURE_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
SLIDE_BG    = RGBColor(0xF3, 0xF5, 0xF8)
SECTION_BG  = RGBColor(0xF0, 0xF2, 0xF5)
CARD_BG     = RGBColor(0xFF, 0xFF, 0xFF)
SUCCESS     = RGBColor(0x02, 0x7A, 0x48)
SUCCESS_LT  = RGBColor(0xEC, 0xFD, 0xF3)
WARNING     = RGBColor(0xB5, 0x47, 0x08)
WARNING_LT  = RGBColor(0xFE, 0xF0, 0xC7)
ERROR       = RGBColor(0xB3, 0x23, 0x18)
ERROR_LT    = RGBColor(0xFE, 0xE4, 0xE2)
ACCENT_DARK = RGBColor(0x1A, 0x2E, 0x5A)

SKILL_DIR = os.path.expanduser("~/.claude/skills/waterplan-pptx")
ASSETS = os.path.join(SKILL_DIR, "assets")

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]

# ─── Helpers ───────────────────────────────────────────────────────

def set_slide_bg(slide, color=SLIDE_BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_accent_bars(slide):
    # Top bar
    top = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Pt(4))
    top.fill.solid(); top.fill.fore_color.rgb = FIREFLY_500; top.line.fill.background()
    # Bottom bar
    bot = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(7.5)-Pt(4), Inches(13.333), Pt(4))
    bot.fill.solid(); bot.fill.fore_color.rgb = FIREFLY_500; bot.line.fill.background()

def add_logo(slide):
    logo_path = os.path.join(ASSETS, "waterplan-logo-dark.png")
    if os.path.exists(logo_path):
        slide.shapes.add_picture(logo_path, Inches(11.2), Inches(0.38), Inches(1.6), Inches(0.32))

def add_slide_number(slide, num, total):
    tb = slide.shapes.add_textbox(Inches(12.4), Inches(7.0), Inches(0.8), Inches(0.3))
    tf = tb.text_frame
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.RIGHT
    r = p.add_run(); r.text = f"{num}/{total}"; r.font.size = Pt(9)
    r.font.color.rgb = GRAY_500; r.font.name = "Inter"

def add_title(slide, text, sub=None, y=0.35):
    tb = slide.shapes.add_textbox(Inches(0.7), Inches(y), Inches(10), Inches(0.7))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = text; r.font.size = Pt(26); r.font.bold = True
    r.font.color.rgb = FIREFLY_700; r.font.name = "Inter"
    if sub:
        tb2 = slide.shapes.add_textbox(Inches(0.7), Inches(y+0.55), Inches(10), Inches(0.4))
        tf2 = tb2.text_frame; p2 = tf2.paragraphs[0]
        r2 = p2.add_run(); r2.text = sub; r2.font.size = Pt(14)
        r2.font.color.rgb = GRAY_700; r2.font.name = "Inter"

def add_card(slide, left, top, width, height, fill=CARD_BG, border_color=None):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if border_color:
        s.line.color.rgb = border_color; s.line.width = Pt(1)
    else:
        s.line.fill.background()
    s.adjustments[0] = 0.06
    return s

def add_section_band(slide, left, top, width, height):
    return add_card(slide, left, top, width, height, fill=SECTION_BG)

def add_pill(slide, left, top, width, height, fill, text="", text_color=WHITE, font_size=11, bold=True):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    s.fill.solid(); s.fill.fore_color.rgb = fill; s.line.fill.background()
    s.adjustments[0] = 0.5
    if text:
        tf = s.text_frame; tf.word_wrap = False
        tf.margin_left = Pt(6); tf.margin_right = Pt(6)
        tf.margin_top = Pt(0); tf.margin_bottom = Pt(0)
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = text; r.font.size = Pt(font_size)
        r.font.color.rgb = text_color; r.font.name = "Inter"; r.font.bold = bold
    return s

def add_text(slide, left, top, width, height, text, size=14, color=GRAY_800, bold=False, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text; r.font.size = Pt(size)
    r.font.color.rgb = color; r.font.name = "Inter"; r.font.bold = bold
    return tf

def add_multiline(slide, left, top, width, height, lines, size=13, color=GRAY_800, line_spacing=1.3, bullet=False):
    """lines = list of (text, bold) or just strings"""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame; tf.word_wrap = True
    for i, line in enumerate(lines):
        if isinstance(line, tuple):
            txt, bld = line
        else:
            txt, bld = line, False
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_after = Pt(4)
        if bullet:
            p.level = 0
            txt = "  " + txt
        r = p.add_run()
        r.text = ("" if not bullet else "  ") + txt
        r.font.size = Pt(size); r.font.color.rgb = color
        r.font.name = "Inter"; r.font.bold = bld
    return tf

def base_slide(num, total):
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide)
    add_accent_bars(slide)
    add_logo(slide)
    add_slide_number(slide, num, total)
    return slide

TOTAL = 12

# ═══════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE SLIDE (dark bg)
# ═══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
set_slide_bg(s, FIREFLY_900)
# Top accent bar
top = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Pt(5))
top.fill.solid(); top.fill.fore_color.rgb = FIREFLY_500; top.line.fill.background()
# Bottom accent bar
bot = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(7.5)-Pt(5), Inches(13.333), Pt(5))
bot.fill.solid(); bot.fill.fore_color.rgb = FIREFLY_500; bot.line.fill.background()

# Logo white
logo_w = os.path.join(ASSETS, "waterplan-logo-white.png")
if os.path.exists(logo_w):
    s.shapes.add_picture(logo_w, Inches(0.7), Inches(0.6), Inches(2.0), Inches(0.4))

# Main title
add_text(s, Inches(0.7), Inches(2.2), Inches(10), Inches(1.2),
         "E-Leads Workshops", size=42, color=PURE_WHITE, bold=True)
add_text(s, Inches(0.7), Inches(3.2), Inches(10), Inches(0.8),
         "Findings, Barriers & Recommendations for Target Tracking Adoption",
         size=22, color=FIREFLY_100)

# Meta info
meta_lines = [
    "3 workshops  |  5 regions  |  8 e-leads  |  April 30 - May 8, 2026",
    "",
    "Regions: FBN (France, Benelux, Nordics)  ·  Iberia  ·  GB + Germany",
]
tb = s.shapes.add_textbox(Inches(0.7), Inches(4.6), Inches(8), Inches(1.2))
tf = tb.text_frame; tf.word_wrap = True
for i, line in enumerate(meta_lines):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    r = p.add_run(); r.text = line; r.font.size = Pt(14)
    r.font.color.rgb = FIREFLY_100; r.font.name = "Inter"

# Prepared for
add_text(s, Inches(0.7), Inches(6.2), Inches(5), Inches(0.4),
         "Prepared for Tobias & Elena  |  Waterplan x CCEP", size=12, color=GRAY_500)

# YC
yc_path = os.path.join(ASSETS, "yc-logo-white.png")
if os.path.exists(yc_path):
    s.shapes.add_picture(yc_path, Inches(11.0), Inches(6.5), Inches(1.5), Inches(0.3))

# ═══════════════════════════════════════════════════════════════════
# SLIDE 2 — EXECUTIVE SUMMARY
# ═══════════════════════════════════════════════════════════════════
s = base_slide(2, TOTAL)
add_title(s, "Executive Summary", sub="What we heard across all 3 workshops")

# Left column — big insight
add_section_band(s, Inches(0.7), Inches(1.5), Inches(5.8), Inches(5.3))
add_text(s, Inches(1.0), Inches(1.7), Inches(5.2), Inches(0.4),
         "THE MAIN FINDING", size=11, color=FIREFLY_500, bold=True)

insight_tb = s.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(5.2), Inches(1.6))
itf = insight_tb.text_frame; itf.word_wrap = True
p = itf.paragraphs[0]
r = p.add_run(); r.text = "E-leads see clear value in Target Tracking, "
r.font.size = Pt(16); r.font.color.rgb = GRAY_900; r.font.name = "Inter"
r = p.add_run(); r.text = "but they won't adopt it "
r.font.size = Pt(16); r.font.color.rgb = GRAY_900; r.font.name = "Inter"; r.font.bold = True
r = p.add_run(); r.text = "until the relationship with Integrum is clear and double work is eliminated."
r.font.size = Pt(16); r.font.color.rgb = GRAY_900; r.font.name = "Inter"

bullets = [
    ("Interest level: Moderate-high, but with conditions", True),
    ("Most used features: 2030 targets, scenario modeling", False),
    ("Least used: Risk maps, future projections, dashboards", False),
    ("Frequency: Not daily. Occasional / need-based only", False),
    ("Key quote: \"Rather than being another tool,\nit should be THE tool\" — Rory, GB", False),
]
tb = s.shapes.add_textbox(Inches(1.0), Inches(3.9), Inches(5.2), Inches(2.8))
tf = tb.text_frame; tf.word_wrap = True
for i, (txt, bld) in enumerate(bullets):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.space_before = Pt(6); p.space_after = Pt(2)
    r = p.add_run(); r.text = "•  " + txt; r.font.size = Pt(13)
    r.font.color.rgb = GRAY_800; r.font.name = "Inter"; r.font.bold = bld

# Right column — 4 KPI cards
cards_data = [
    ("5 REGIONS", "FBN, Iberia, GB, Germany, Nordics", FIREFLY_500),
    ("8 E-LEADS", "All active environmental managers", FIREFLY_700),
    ("18 THEMES", "Identified from feedback analysis", FIREFLY_500),
    ("6 DECISIONS", "Required from CCP global", WARNING),
]
for i, (title, desc, color) in enumerate(cards_data):
    cy = Inches(1.5) + Inches(i * 1.3)
    c = add_card(s, Inches(6.8), cy, Inches(5.7), Inches(1.1))
    add_text(s, Inches(7.1), cy + Inches(0.15), Inches(2.5), Inches(0.35),
             title, size=18, color=color, bold=True)
    add_text(s, Inches(7.1), cy + Inches(0.55), Inches(5.0), Inches(0.35),
             desc, size=12, color=GRAY_700)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 3 — TOP 5 TAKEAWAYS (page 1: #1-3)
# ═══════════════════════════════════════════════════════════════════
s = base_slide(3, TOTAL)
add_title(s, "Top 5 Takeaways for the Global Meeting", sub="Critical insights from all workshops — Part 1")

takeaways_1 = [
    {
        "num": "1",
        "title": "Waterplan is seen as \"yet another tool\"",
        "evidence": "Jos (Benelux): \"We have 4 project tracking tools. Sites are complaining.\" Rory (GB): wants it to be \"the tool, not another tool.\"",
        "action": "Define Waterplan's role vs Integrum. Without this, adoption stays marginal.",
        "color": ERROR,
        "bg": ERROR_LT,
    },
    {
        "num": "2",
        "title": "Integrum sync is manual, weekly, one-way — blocks trust",
        "evidence": "James (Nordics): \"Not a real-time process.\" Antonio (Iberia): \"Changes here don't go back to Integrum.\" Laurie: \"I won't update a 3rd tool.\"",
        "action": "Plan bidirectional integration or define Waterplan as the single source for water.",
        "color": ERROR,
        "bg": ERROR_LT,
    },
    {
        "num": "3",
        "title": "E-leads need energy + carbon, not just water",
        "evidence": "Mentioned in ALL 3 workshops. Laurie: \"I wouldn't want one platform for water and another for climate.\" Annual target Excel includes both.",
        "action": "Evaluate if Target Tracking scope should expand to energy/carbon.",
        "color": WARNING,
        "bg": WARNING_LT,
    },
]

for i, t in enumerate(takeaways_1):
    cy = Inches(1.55) + Inches(i * 1.85)
    # Number circle
    add_pill(s, Inches(0.7), cy + Inches(0.1), Inches(0.45), Inches(0.45),
             t["color"], t["num"], PURE_WHITE, 18)
    # Card
    add_card(s, Inches(1.35), cy, Inches(11.3), Inches(1.65), fill=t["bg"], border_color=t["color"])
    # Title
    add_text(s, Inches(1.55), cy + Inches(0.1), Inches(10.8), Inches(0.35),
             t["title"], size=16, color=GRAY_900, bold=True)
    # Evidence
    add_text(s, Inches(1.55), cy + Inches(0.5), Inches(5.4), Inches(0.9),
             "Evidence: " + t["evidence"], size=11, color=GRAY_700)
    # Action
    add_pill(s, Inches(7.2), cy + Inches(0.5), Inches(1.2), Inches(0.28),
             t["color"], "ACTION", PURE_WHITE, 9)
    add_text(s, Inches(7.2), cy + Inches(0.85), Inches(5.2), Inches(0.65),
             t["action"], size=12, color=GRAY_900, bold=True)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 4 — TOP 5 TAKEAWAYS (page 2: #4-5)
# ═══════════════════════════════════════════════════════════════════
s = base_slide(4, TOTAL)
add_title(s, "Top 5 Takeaways for the Global Meeting", sub="Critical insights from all workshops — Part 2")

takeaways_2 = [
    {
        "num": "4",
        "title": "Navigation is not intuitive. Naming is not CCP language.",
        "evidence": "Antonio (Iberia): \"'Responder' to reach targets? Not intuitive.\" Manuel (Germany): \"Tough to explain to site managers what you can see.\" Monika (GB): \"Would site managers explore? Possibly not.\"",
        "action": "Ship homepage with shortcuts + rename modules to CCP terminology. Already in progress with Elena.",
        "color": WARNING,
        "bg": WARNING_LT,
    },
    {
        "num": "5",
        "title": "2030 targets can be edited by anyone — governance risk",
        "evidence": "Manuel (Germany): \"We told site managers: please don't change targets. But there's no lock.\" Antonio (Iberia): \"Can I accidentally duplicate a project into Germany?\" — Answer: yes.",
        "action": "Implement target lock + role-based permissions immediately. Define policy with CCP global.",
        "color": ERROR,
        "bg": ERROR_LT,
    },
]

for i, t in enumerate(takeaways_2):
    cy = Inches(1.55) + Inches(i * 2.4)
    add_pill(s, Inches(0.7), cy + Inches(0.1), Inches(0.45), Inches(0.45),
             t["color"], t["num"], PURE_WHITE, 18)
    add_card(s, Inches(1.35), cy, Inches(11.3), Inches(2.1), fill=t["bg"], border_color=t["color"])
    add_text(s, Inches(1.55), cy + Inches(0.12), Inches(10.8), Inches(0.35),
             t["title"], size=16, color=GRAY_900, bold=True)
    add_text(s, Inches(1.55), cy + Inches(0.55), Inches(5.4), Inches(1.2),
             "Evidence: " + t["evidence"], size=11, color=GRAY_700)
    add_pill(s, Inches(7.2), cy + Inches(0.55), Inches(1.2), Inches(0.28),
             t["color"], "ACTION", PURE_WHITE, 9)
    add_text(s, Inches(7.2), cy + Inches(0.92), Inches(5.2), Inches(0.9),
             t["action"], size=12, color=GRAY_900, bold=True)

# Bottom insight box
add_card(s, Inches(0.7), Inches(5.7), Inches(11.95), Inches(1.0), fill=FIREFLY_50, border_color=FIREFLY_500)
add_text(s, Inches(1.0), Inches(5.85), Inches(11.4), Inches(0.7),
         "KEY PATTERN: The top 3 barriers (Integrum duplication, energy/carbon scope, UX) appeared in ALL 3 workshops across ALL 5 regions. This is not a regional issue — it is a structural blocker that requires a global decision.",
         size=13, color=FIREFLY_900, bold=True)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 5 — ADOPTION BARRIERS (visual)
# ═══════════════════════════════════════════════════════════════════
s = base_slide(5, TOTAL)
add_title(s, "Why E-Leads Are Not Adopting Yet", sub="6 barriers identified — ordered by frequency and impact")

barriers = [
    ("Double work with Integrum", "All regions", "Integrum = source of truth. Waterplan = extra effort.\nNo bidirectional sync. E-leads update Integrum only.", ERROR, "CRITICAL"),
    ("Water-only scope", "All regions", "E-leads manage water + energy + carbon together.\nWaterplan only covers water → always partial.", WARNING, "HIGH"),
    ("UX & naming friction", "4 of 5 regions", "Non-intuitive navigation. Too much info for site managers.\nModule names don't match CCP vocabulary.", WARNING, "HIGH"),
    ("No permission controls", "Germany, Iberia", "Anyone can edit 2030 targets or other BU's projects.\nNo lock, no audit trail, no role separation.", ERROR, "HIGH"),
    ("No global mandate", "All regions", "E-leads treat Waterplan as optional.\nNo one told them they must use it.", WARNING, "MEDIUM"),
    ("Tool fatigue at sites", "Benelux", "4 project-tracking tools + 7 risk assessments.\nSites refuse to adopt yet another platform.", WARNING, "MEDIUM"),
]

for i, (title, who, desc, color, sev) in enumerate(barriers):
    row = i // 2
    col = i % 2
    cx = Inches(0.7) + Inches(col * 6.3)
    cy = Inches(1.5) + Inches(row * 1.85)

    add_card(s, cx, cy, Inches(5.95), Inches(1.65))

    # Severity pill
    add_pill(s, cx + Inches(0.15), cy + Inches(0.15), Inches(1.0), Inches(0.28),
             color, sev, PURE_WHITE, 9)

    # Who pill
    add_pill(s, cx + Inches(1.25), cy + Inches(0.15), Inches(1.4), Inches(0.28),
             GRAY_50, who, GRAY_700, 9, bold=False)

    # Title
    add_text(s, cx + Inches(0.2), cy + Inches(0.55), Inches(5.5), Inches(0.3),
             title, size=14, color=GRAY_900, bold=True)

    # Description
    add_text(s, cx + Inches(0.2), cy + Inches(0.9), Inches(5.5), Inches(0.6),
             desc, size=11, color=GRAY_700)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 6 — FEEDBACK THEMES TABLE
# ═══════════════════════════════════════════════════════════════════
s = base_slide(6, TOTAL)
add_title(s, "Consolidated Feedback Themes", sub="18 themes identified — top 10 shown by priority")

# Table
headers = ["#", "Theme", "Regions", "Priority", "Action Type", "Our Recommendation"]
rows = [
    headers,
    ["1", "Integrum duplication", "All (5/5)", "HIGH", "Data + CCP Global", "Bidirectional sync or single source"],
    ["2", "Permissions & lock targets", "DE, Iberia", "HIGH", "Product", "Lock targets + role matrix now"],
    ["3", "Navigation & naming", "4 of 5", "HIGH", "Product", "Homepage + CCP naming (in progress)"],
    ["4", "Yearly target setting", "All (5/5)", "HIGH", "Product + CCP", "Expand feature, align with 2030 logic"],
    ["5", "Comments & activity log", "FBN, DE", "HIGH", "Product", "Add to mitigation actions module"],
    ["6", "Energy/carbon scope", "All (5/5)", "MEDIUM", "Product + CCP", "Evaluate expansion with CCP global"],
    ["7", "Tool saturation at sites", "Benelux", "MEDIUM", "CCP Global", "Consolidation plan needed"],
    ["8", "Sites hide projects", "Benelux", "MEDIUM", "CS + CCP", "Scenarios as 'safe space' to explore"],
    ["9", "Regional dashboard", "GB, Iberia", "MEDIUM", "Product", "At-a-glance view per BU"],
    ["10", "Financial view by year", "Iberia", "MEDIUM", "Product", "Capex/opex timeline in scenarios"],
]

n_rows = len(rows); n_cols = len(rows[0])
tbl_shape = s.shapes.add_table(n_rows, n_cols,
    Inches(0.7), Inches(1.5), Inches(11.95), Inches(5.2))
tbl = tbl_shape.table

# Column widths
col_widths = [Inches(0.4), Inches(2.5), Inches(1.3), Inches(1.0), Inches(2.5), Inches(4.25)]
for ci, w in enumerate(col_widths):
    tbl.columns[ci].width = w

priority_colors = {"HIGH": ERROR, "MEDIUM": WARNING, "LOW": SUCCESS}

for ri, row in enumerate(rows):
    for ci, val in enumerate(row):
        cell = tbl.cell(ri, ci)
        cell.text = ""
        tf = cell.text_frame; tf.word_wrap = True
        tf.margin_left = Pt(4); tf.margin_right = Pt(4)
        tf.margin_top = Pt(3); tf.margin_bottom = Pt(3)
        p = tf.paragraphs[0]
        r = p.add_run(); r.text = val; r.font.name = "Inter"

        if ri == 0:
            cell.fill.solid(); cell.fill.fore_color.rgb = FIREFLY_700
            r.font.size = Pt(10); r.font.color.rgb = PURE_WHITE; r.font.bold = True
        else:
            cell.fill.solid()
            cell.fill.fore_color.rgb = PURE_WHITE if ri % 2 == 1 else FIREFLY_50
            r.font.size = Pt(10); r.font.color.rgb = GRAY_900
            if ci == 0:
                r.font.bold = True; r.font.color.rgb = FIREFLY_500
            if ci == 3 and val in priority_colors:
                r.font.bold = True; r.font.color.rgb = priority_colors[val]

# ═══════════════════════════════════════════════════════════════════
# SLIDE 7 — WHAT E-LEADS SAID (quotes)
# ═══════════════════════════════════════════════════════════════════
s = base_slide(7, TOTAL)
add_title(s, "In Their Own Words", sub="Direct quotes from the workshops")

quotes = [
    {
        "text": "\"We have four project tracking tools.\nThe technical departments are really complaining.\"",
        "who": "Jos Peeters — Benelux",
        "tag": "Tool fatigue",
    },
    {
        "text": "\"Rather than being another tool, it should be THE tool.\nI think it's got potential.\"",
        "who": "Rory Goodwin — GB",
        "tag": "Aspiration",
    },
    {
        "text": "\"I really like the project and target tracking part.\nIt's much more intuitive than doing everything in Excel.\"",
        "who": "Manuel Schlapka — Germany",
        "tag": "Value seen",
    },
    {
        "text": "\"The tool is very practical. But if Integrum and Waterplan\nare not connected, I prefer to work only once.\"",
        "who": "Antonio Agudo — Iberia",
        "tag": "Condition",
    },
    {
        "text": "\"I wouldn't want to go on one platform for water\nand another for climate projects.\"",
        "who": "Laurie Sarkissian — France",
        "tag": "Scope gap",
    },
    {
        "text": "\"We told site managers: please don't change any target.\nBut there is no lock. I just hope for no changes.\"",
        "who": "Manuel Schlapka — Germany",
        "tag": "Governance",
    },
]

for i, q in enumerate(quotes):
    col = i % 2
    row = i // 2
    cx = Inches(0.7) + Inches(col * 6.3)
    cy = Inches(1.5) + Inches(row * 1.85)

    add_card(s, cx, cy, Inches(5.95), Inches(1.65))

    # Tag pill
    tag_colors = {
        "Tool fatigue": ERROR, "Aspiration": SUCCESS, "Value seen": SUCCESS,
        "Condition": WARNING, "Scope gap": WARNING, "Governance": ERROR,
    }
    tc = tag_colors.get(q["tag"], GRAY_500)
    add_pill(s, cx + Inches(0.2), cy + Inches(0.15), Inches(1.4), Inches(0.25),
             tc, q["tag"], PURE_WHITE, 9)

    # Quote
    add_text(s, cx + Inches(0.2), cy + Inches(0.5), Inches(5.5), Inches(0.8),
             q["text"], size=12, color=GRAY_900)

    # Attribution
    add_text(s, cx + Inches(0.2), cy + Inches(1.3), Inches(5.5), Inches(0.25),
             "— " + q["who"], size=10, color=FIREFLY_500, bold=True)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 8 — ACTION PLAN (Short term)
# ═══════════════════════════════════════════════════════════════════
s = base_slide(8, TOTAL)
add_title(s, "Recommended Action Plan", sub="What we should do — and when")

# 3 time horizons as columns
horizons = [
    {
        "label": "SHORT TERM", "period": "Next 2-4 weeks",
        "color": SUCCESS, "bg": SUCCESS_LT,
        "actions": [
            "Fix production volume bug (Germany)",
            "Implement lock for 2030 targets",
            "Send workshop summary to e-leads",
            "Run parallel pilot with Iberia\n(annual update in both platforms)",
            "Ship homepage with shortcuts",
            "Apply CCP naming to navigation",
        ],
    },
    {
        "label": "MEDIUM TERM", "period": "Next 2-3 months",
        "color": WARNING, "bg": WARNING_LT,
        "actions": [
            "Define permission policy per role\nwith CCP global",
            "Add comments + activity log\nto mitigation actions",
            "Create regional dashboard for e-leads",
            "Complete yearly targets feature\naligned with 2030 logic",
            "Train site managers (Iberia, Germany)",
            "Add financial view (capex by year)",
        ],
    },
    {
        "label": "LONG TERM", "period": "Next 6 months",
        "color": FIREFLY_700, "bg": FIREFLY_50,
        "actions": [
            "Build bidirectional Integrum\nintegration roadmap",
            "Evaluate energy/carbon expansion\nin Target Tracking",
            "Create global adoption playbook",
            "Consolidate Excel tools\n(TCW, Top 10, annual targets)",
            "Position Waterplan as primary\nwater platform for CCEP",
        ],
    },
]

for hi, h in enumerate(horizons):
    cx = Inches(0.7) + Inches(hi * 4.15)
    # Header pill
    add_pill(s, cx, Inches(1.5), Inches(3.85), Inches(0.4),
             h["color"], f"{h['label']}  |  {h['period']}", PURE_WHITE, 12)
    # Card
    add_card(s, cx, Inches(2.05), Inches(3.85), Inches(4.85), fill=h["bg"])
    # Actions
    for ai, action in enumerate(h["actions"]):
        ay = Inches(2.2) + Inches(ai * 0.78)
        # Small white card per action
        add_card(s, cx + Inches(0.1), ay, Inches(3.65), Inches(0.68))
        add_text(s, cx + Inches(0.25), ay + Inches(0.08), Inches(3.4), Inches(0.55),
                 action, size=11, color=GRAY_900)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 9 — IDEAS TO DISCUSS
# ═══════════════════════════════════════════════════════════════════
s = base_slide(9, TOTAL)
add_title(s, "Ideas to Discuss Together", sub="Proposals for the global meeting")

ideas = [
    ("Parallel validation period", "Sites keep reporting in Integrum. Waterplan runs in parallel.\nIf data matches → Waterplan becomes trusted source.", "Iberia is ready for this pilot."),
    ("Scenarios as 'safe space'", "Let sites explore projects in Waterplan without triggering\nofficial reporting. Reduces fear of uploading.", "Solves the 'hidden opex projects' problem."),
    ("Regional pilots first", "Start with Iberia (Antonio) + Germany (Manuel) — most\nadvanced users. Prove value before global rollout.", "Faster learning, lower risk."),
    ("Define minimum criteria per site", "Site needs: projects loaded, 2030 target locked,\n1 scenario modeled, 1 trained manager.", "Clear checklist for readiness."),
    ("Global adoption playbook", "3 levels: L1 = View only, L2 = Operational,\nL3 = Reporting from Waterplan.", "Scalable, not all-or-nothing."),
    ("Migrate key Excels first", "Annual target setting, True Cost of Water, Top 10\npractices → move to Waterplan.", "Quick wins that reduce tool fatigue."),
    ("Define reporting outputs", "Which reports should come from Waterplan?\nFor whom? Replace existing or complement?", "Without this, no pull from stakeholders."),
    ("Capex linearization view", "Show investment needed per year to meet 2030.\nHelps e-leads plan budgets.", "Antonio requested this specifically."),
]

for i, (title, desc, note) in enumerate(ideas):
    col = i % 2
    row = i // 2
    cx = Inches(0.7) + Inches(col * 6.3)
    cy = Inches(1.4) + Inches(row * 1.43)

    add_card(s, cx, cy, Inches(5.95), Inches(1.28))

    # Number
    add_pill(s, cx + Inches(0.15), cy + Inches(0.12), Inches(0.32), Inches(0.32),
             FIREFLY_500, str(i+1), PURE_WHITE, 11)

    add_text(s, cx + Inches(0.6), cy + Inches(0.1), Inches(5.1), Inches(0.28),
             title, size=13, color=GRAY_900, bold=True)
    add_text(s, cx + Inches(0.6), cy + Inches(0.42), Inches(5.1), Inches(0.5),
             desc, size=10, color=GRAY_700)
    add_text(s, cx + Inches(0.6), cy + Inches(0.95), Inches(5.1), Inches(0.25),
             "→ " + note, size=10, color=FIREFLY_500, bold=True)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 10 — KEY QUESTIONS
# ═══════════════════════════════════════════════════════════════════
s = base_slide(10, TOTAL)
add_title(s, "Key Questions for the Global Decision", sub="These need answers before we can move forward")

q_groups = [
    ("PLATFORM ROLE", FIREFLY_700, [
        "What is Waterplan's role vs Integrum for water?",
        "Complementary or replacement?",
        "Who decides on Integrum changes?",
    ]),
    ("GOVERNANCE", ERROR, [
        "Who owns 2030 targets — e-leads or central?",
        "Should site managers edit or view-only?",
        "Do we need a formal approval flow?",
    ]),
    ("ADOPTION", SUCCESS, [
        "Is Waterplan mandatory or optional for e-leads?",
        "Should we start with pilots or go global?",
        "What does 'success' look like in 6 months?",
    ]),
    ("SCOPE", WARNING, [
        "Should Waterplan expand to energy/carbon?",
        "Which reports should come from Waterplan?",
        "Is there a tool rationalization plan at CCP?",
    ]),
]

for gi, (label, color, questions) in enumerate(q_groups):
    col = gi % 2
    row = gi // 2
    cx = Inches(0.7) + Inches(col * 6.3)
    cy = Inches(1.5) + Inches(row * 2.7)

    add_section_band(s, cx, cy, Inches(5.95), Inches(2.4))
    add_pill(s, cx + Inches(0.2), cy + Inches(0.2), Inches(2.2), Inches(0.32),
             color, label, PURE_WHITE, 11)

    for qi, q in enumerate(questions):
        qy = cy + Inches(0.7) + Inches(qi * 0.5)
        add_text(s, cx + Inches(0.4), qy, Inches(5.3), Inches(0.4),
                 f"{qi+1}.  {q}", size=13, color=GRAY_900)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 11 — RISKS IF WE DON'T ACT
# ═══════════════════════════════════════════════════════════════════
s = base_slide(11, TOTAL)
add_title(s, "Risks If We Don't Decide", sub="What happens if we leave things as they are")

risks = [
    ("Irrelevance", "Without a global mandate, each e-lead decides alone.\nMost will choose Integrum + their own Excels.\nWaterplan becomes a nice-to-have nobody opens."),
    ("Tool fatigue rebellion", "Sites already complain about 4+ tracking tools.\nAdding a 5th without removing one = active resistance.\nJos explicitly flagged this during an ISO 46001 audit."),
    ("Data divergence", "Without bidirectional sync, Integrum and Waterplan\nwill show different numbers. Both lose credibility.\nReporting becomes unreliable."),
    ("Governance incident", "Without locked targets and role permissions,\nit's a matter of time before someone edits a 2030\ntarget by mistake. Recovery = painful."),
]

for i, (title, desc) in enumerate(risks):
    col = i % 2
    row = i // 2
    cx = Inches(0.7) + Inches(col * 6.3)
    cy = Inches(1.5) + Inches(row * 2.7)

    add_card(s, cx, cy, Inches(5.95), Inches(2.4), fill=PURE_WHITE, border_color=ERROR)

    # Warning icon (text-based)
    add_pill(s, cx + Inches(0.2), cy + Inches(0.2), Inches(0.35), Inches(0.35),
             ERROR, "!", PURE_WHITE, 16)

    add_text(s, cx + Inches(0.7), cy + Inches(0.2), Inches(5.0), Inches(0.3),
             title, size=16, color=ERROR, bold=True)

    add_text(s, cx + Inches(0.3), cy + Inches(0.7), Inches(5.4), Inches(1.5),
             desc, size=13, color=GRAY_800)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 12 — NEXT STEPS / CLOSE
# ═══════════════════════════════════════════════════════════════════
s = base_slide(12, TOTAL)
add_title(s, "Proposed Next Steps")

steps = [
    ("THIS WEEK", FIREFLY_500, [
        "Align with Tobias on Waterplan vs Integrum positioning",
        "Confirm permission policy (who edits, who views)",
        "Deploy target lock feature",
    ]),
    ("THIS MONTH", FIREFLY_700, [
        "Launch parallel pilot with Iberia (annual update)",
        "Ship homepage + CCP naming changes",
        "Send follow-up to all e-leads with action plan",
    ]),
    ("THIS QUARTER", FIREFLY_900, [
        "Define bidirectional integration roadmap with CCP IT",
        "Deliver yearly targets feature aligned with 2030",
        "Evaluate energy/carbon scope decision",
        "Train site managers in pilot regions",
    ]),
]

for si, (label, color, items) in enumerate(steps):
    cx = Inches(0.7) + Inches(si * 4.15)
    add_pill(s, cx, Inches(1.5), Inches(3.85), Inches(0.4), color, label, PURE_WHITE, 13)

    add_card(s, cx, Inches(2.1), Inches(3.85), Inches(3.8))

    for ii, item in enumerate(items):
        iy = Inches(2.3) + Inches(ii * 0.85)
        # Checkmark circle
        add_pill(s, cx + Inches(0.2), iy + Inches(0.05), Inches(0.28), Inches(0.28),
                 FIREFLY_100, "○", color, 12, bold=False)
        add_text(s, cx + Inches(0.6), iy, Inches(3.0), Inches(0.7),
                 item, size=12, color=GRAY_900)

# Bottom CTA
add_card(s, Inches(0.7), Inches(6.15), Inches(11.95), Inches(0.8), fill=FIREFLY_700)
add_text(s, Inches(1.0), Inches(6.25), Inches(11.4), Inches(0.55),
         "The e-leads are ready. They see the value. What they need now is a clear signal from global leadership that Target Tracking is THE path forward.",
         size=15, color=PURE_WHITE, bold=True, anchor=MSO_ANCHOR.MIDDLE)


# ═══════════════════════════════════════════════════════════════════
# SAVE
# ═══════════════════════════════════════════════════════════════════
output_dir = os.path.expanduser("~/Downloads")
output_path = os.path.join(output_dir, "CCEP-ELeads-Workshop-Findings.pptx")
prs.save(output_path)
print(f"Saved to {output_path}")
